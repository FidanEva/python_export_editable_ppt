from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION, XL_DATA_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from services.excel_parser import get_sentiment_counts, get_company_sentiment_counts
import logging
import pandas as pd
import os
from PIL import Image

logger = logging.getLogger(__name__)

# Define colors
SLIDE_BG_COLOR = RGBColor(245, 245, 245)  # Light gray
CHART_BG_COLOR = RGBColor(255, 255, 255)  # White
DEFAULT_COLOR = RGBColor(214, 55, 64)    # Red
HEADER_HEIGHT = 0.8

# Define sentiment colors
SENTIMENT_COLORS = {
    1: RGBColor(69, 194, 126),    # Positive - Green
    0: RGBColor(255, 191, 0),     # Neutral - Yellow
    -1: RGBColor(246, 1, 64)      # Negative - Red
}

# Define metric icons
METRIC_ICONS = {
    'Post sayı': '🔔',
    'Şərh sayı': '💬',
    'Bəyənmə sayı': '👍',
    'Paylaşım sayı': '📤',
    'Baxış sayı': '👁'
}

CHARTS_ICONS = {
    'Sentiment Trend': '📉',
    'Sentiment Distribution': '📊',
    "Time Distribution": '🕒',
}
def format_number_with_k(value):
    """Format number with k for thousands"""
    if value >= 1000:
        return f"{value/1000:.1f}k".replace(".0k", "k")
    return f"{value:,}"

def format_chart_axes(chart, category_font_size=6, value_font_size=6, rotation_angle=-45):
    """Helper function to consistently format chart axes"""
    if hasattr(chart, 'category_axis'):
        chart.category_axis.tick_labels.font.size = Pt(category_font_size)
        chart.category_axis.tick_labels.font.bold = False
        if hasattr(chart.category_axis, 'axis_title') and chart.category_axis.axis_title:
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(category_font_size)
        # Set labels at specified angle
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        chart.category_axis.tick_label_rotation = rotation_angle
        chart.category_axis.tick_labels.font.color.rgb = RGBColor(89, 89, 89)
        chart.category_axis.has_major_gridlines = False
    
    if hasattr(chart, 'value_axis'):
        chart.value_axis.tick_labels.font.size = Pt(value_font_size)
        chart.value_axis.tick_labels.font.bold = False
        if hasattr(chart.value_axis, 'axis_title') and chart.value_axis.axis_title:
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(value_font_size)
        chart.value_axis.tick_labels.font.color.rgb = RGBColor(89, 89, 89)
        chart.value_axis.has_major_gridlines = False
def add_slide_header(slide, company_logo_path, start_date, end_date, title, title_color=DEFAULT_COLOR, template_color=DEFAULT_COLOR):
    """Helper function to add a consistent header to slides"""

    # Set constants
    FULL_SLIDE_WIDTH_INCHES = 13.33
    HEADER_WIDTH_RATIO = 0.95
    HEADER_WIDTH = FULL_SLIDE_WIDTH_INCHES * HEADER_WIDTH_RATIO
    HEADER_LEFT_MARGIN = (FULL_SLIDE_WIDTH_INCHES - HEADER_WIDTH) / 2

    # Add header background with SLIDE_BG_COLOR
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(HEADER_LEFT_MARGIN), Inches(0),
        Inches(HEADER_WIDTH), Inches(HEADER_HEIGHT)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = SLIDE_BG_COLOR
    header.line.fill.background()  # Remove border

    # Add company logo (left-aligned within the header)
    if os.path.exists(company_logo_path):
        slide.shapes.add_picture(
            company_logo_path,
            Inches(HEADER_LEFT_MARGIN + 0.2), Inches(0.2),
            height=Inches(0.4)
        )

    # Add title in the center of the header
    title_box = slide.shapes.add_textbox(
        Inches(HEADER_LEFT_MARGIN + (HEADER_WIDTH - 6) / 2), Inches(0.2),
        Inches(6), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"📊 {title}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = title_color

    # Add date on the far right, a bit lower
    date_box = slide.shapes.add_textbox(
        Inches(HEADER_LEFT_MARGIN + HEADER_WIDTH - 3), Inches(0.35),
        Inches(2.8), Inches(0.6)
    )
    tf = date_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"📅 {start_date} - {end_date}"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.RIGHT
    p.font.color.rgb = title_color

    # Add divider line under the header, same width
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(HEADER_LEFT_MARGIN), Inches(HEADER_HEIGHT),
        Inches(HEADER_WIDTH), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = template_color
    line.line.fill.background()
    
    line.shadow.inherit = False

def add_side_line(slide, template_color=DEFAULT_COLOR):
    """Helper function to add side line with only right corners rounded"""
    # Add side line (70% of slide height, centered vertically)
    slide_height = Inches(7.5)
    line_height = slide_height * 0.6
    start_y = (slide_height - line_height) / 2

    # Create a rounded rectangle for the side line
    side_line = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0),
        start_y,
        Inches(0.2),
        line_height
    )
    
    if len(side_line.adjustments) >= 5:
        side_line.adjustments[0] = 1
        side_line.adjustments[1] = 0
        side_line.adjustments[2] = 1
        side_line.adjustments[3] = 1
        side_line.adjustments[4] = 0
    
    # Apply fill and remove border
    side_line.fill.solid()
    side_line.fill.fore_color.rgb = template_color
    side_line.line.fill.background()
    
# Function to apply consistent chart formatting
def apply_chart_formatting(chart, use_legend=True, legend_position=XL_LEGEND_POSITION.TOP, 
                         category_font_size=6, value_font_size=6, 
                         title=None, title_size=14, icon='📊', graph_color = DEFAULT_COLOR):
    """Helper function for consistent chart formatting across slides"""
    
    # Basic chart settings
    chart.chart_style = 2  # White background
    chart.plots[0].has_major_gridlines = False
    chart.plots[0].has_minor_gridlines = False
    
    # Legend settings
    chart.has_legend = use_legend
    if use_legend:
        chart.legend.position = legend_position
        chart.legend.font.size = Pt(12)
        
    # Title settings
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = f"{icon} {title}"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(title_size)
        chart.chart_title.text_frame.paragraphs[0].font.bold = False
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = graph_color
        
    # Axis formatting
    if hasattr(chart, 'category_axis'):
        chart.category_axis.tick_labels.font.size = Pt(category_font_size)
        chart.category_axis.tick_labels.font.bold = False
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        chart.category_axis.tick_label_rotation = -45  # 45-degree angle
        chart.category_axis.tick_labels.font.color.rgb = RGBColor(89, 89, 89)
        chart.category_axis.has_major_gridlines = False
        
    if hasattr(chart, 'value_axis'):
        chart.value_axis.tick_labels.font.size = Pt(value_font_size)
        chart.value_axis.tick_labels.font.bold = False
        chart.value_axis.tick_labels.font.color.rgb = RGBColor(89, 89, 89)
        chart.value_axis.has_major_gridlines = False

def apply_sentiment_colors(chart):
    """Helper function to apply consistent sentiment colors to chart series"""
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        # The series order should match [1, 0, -1] for [Positive, Neutral, Negative]
        sentiment_value = [1, 0, -1][i]
        # series.format.fill.fore_color.rgb = SENTIMENT_COLORS[sentiment_value]
        series.format.fill.fore_color.rgb = SENTIMENT_COLORS[list(SENTIMENT_COLORS.keys())[i]]

def add_bg_box(slide, left, top, width, height, color=RGBColor(255, 255, 255)):
    """Helper function to add a background box with rounded corners"""
    bg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg_box.fill.solid()
    bg_box.fill.fore_color.rgb = color
    bg_box.line.fill.background()  # Remove border
    bg_box.shadow.inherit = False
    bg_box.adjustments[0] = 0.01  # Adjust roundness if needed
    return bg_box

def create_sentiment_line_chart(slide, x, y, cx, cy, chart_data, title, icon=CHARTS_ICONS['Sentiment Trend'], graph_color=DEFAULT_COLOR):
    """Helper function to create a sentiment line chart."""
    
    add_bg_box(slide, x, y, cx, cy, color=CHART_BG_COLOR)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.size = Pt(12)

    apply_chart_formatting(chart, title=title, icon=icon, graph_color=graph_color)
    
    for i, series in enumerate(chart.series):
        series.format.line.color.rgb = SENTIMENT_COLORS[list(SENTIMENT_COLORS.keys())[i]]
        series.format.line.width = Pt(2)
        
    return chart

def create_sentiment_donut_chart(slide, x, y, cx, cy, sentiment_counts, graph_color = DEFAULT_COLOR, title = f"{CHARTS_ICONS['Time Distribution']} Bank postlarının sentiment bölgüsü"):
    """Helper function to create a sentiment donut chart."""
    
    add_bg_box(slide, x, y, cx, cy, color=CHART_BG_COLOR)

    donut_data = ChartData()
    donut_data.categories = ['Positive', 'Neutral', 'Negative']
    
    values = [
        int(sentiment_counts.get(1, 0)),
        int(sentiment_counts.get(0, 0)),
        int(sentiment_counts.get(-1, 0))
    ]
    
    donut_data.add_series('', values)

    donut = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, donut_data
    ).chart

    donut.has_legend = True
    donut.legend.position = XL_LEGEND_POSITION.TOP
    donut.legend.font.size = Pt(12)

    donut.has_title = True
    donut.chart_title.text_frame.text = title
    donut.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    donut.chart_title.text_frame.paragraphs[0].font.bold = False
    donut.chart_title.text_frame.paragraphs[0].font.color.rgb = graph_color

    donut.chart_style = 2  # White background

    total = sum(values)

    for i, point in enumerate(donut.series[0].points):
        current_value = values[i]
    
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = SENTIMENT_COLORS[list(SENTIMENT_COLORS.keys())[i]]

        if current_value == 0:
            point.has_data_label = False
            continue

        point.has_data_label = True
        point.data_label.font.bold = True
        
        current_value = values[i]
        percentage = (current_value / total) * 100 if total > 0 else 0
        point.data_label.text_frame.text = f"{current_value:,} ({percentage:.1f}%)"

        for paragraph in point.data_label.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)
    
    return donut

def hex_to_rgbcolor(hex_color):
    if isinstance(hex_color, str) and hex_color.startswith("#") and len(hex_color) == 7:
        r = int(hex_color[1:3], 16)
        g = int(hex_color[3:5], 16)
        b = int(hex_color[5:7], 16)
        return RGBColor(r, g, b)
    return hex_color

def create_ppt(data_frames, output_path, start_date, end_date, company_name, company_logo_path, mediaeye_logo_path, neurotime_logo_path, competitor_logo_paths=None, positive_links=None, negative_links=None, positive_posts=None, negative_posts=None, has_competitors=True, template_color=None, title_color=None, graph_color=None):
    try:
        logger.debug("Creating PowerPoint presentation")
        logger.debug("TEMPLATE COLOR: %s", template_color)
        prs = Presentation()
        
        # Set slide dimensions to landscape (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        

        if template_color is not None:
            template_color = hex_to_rgbcolor(template_color)
        else:
            template_color = DEFAULT_COLOR

        if title_color is not None:
            title_color = hex_to_rgbcolor(title_color)
        else:
            title_color = DEFAULT_COLOR

        if graph_color is not None:
            graph_color = hex_to_rgbcolor(graph_color)
        else:
            graph_color = DEFAULT_COLOR

        
        # region  First slide - Title slide with logos and text
        logger.debug("Creating title slide")
        title_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        # Constants
        SLIDE_WIDTH = 13.33
        LEFT_HALF_CENTER = Inches(SLIDE_WIDTH / 4)        # ~3.33"
        RIGHT_HALF_CENTER = Inches(3 * SLIDE_WIDTH / 4)   # ~10"

        LEFT_CONTENT_WIDTH = Inches(3.5)
        LOGO_WIDTH = Inches(5)
        SPACING_SMALL = Inches(0.1)
        SPACING_MEDIUM = Inches(0.2)

        # Remove default textboxes
        for shape in title_slide.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)

        # Background color
        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR

        # Top offset
        top = Inches(1.0)

        # Calculate centered left margin
        LEFT_MARGIN = LEFT_HALF_CENTER - LEFT_CONTENT_WIDTH / 2

        # NeuroTime logo (top left half)
        if os.path.exists(neurotime_logo_path):
            title_slide.shapes.add_picture(neurotime_logo_path, LEFT_MARGIN, top, width=Inches(4))
            top += Inches(1.5)

        # Main Title
        title_box = title_slide.shapes.add_textbox(LEFT_MARGIN, top, LEFT_CONTENT_WIDTH, Inches(1.2))
        tf = title_box.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "MARKA VARLIĞININ\nMEDİA ÜZƏRİNDƏN\nQİYMƏTLƏNDİRİLMƏSİ"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 123, 191)
        p.line_spacing = Pt(24)
        top += Inches(1.2)

        # Divider Line
        line = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, top, LEFT_CONTENT_WIDTH, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 123, 191)
        line.line.fill.background()
        top += SPACING_SMALL

        # Description
        desc_box = title_slide.shapes.add_textbox(LEFT_MARGIN, top, LEFT_CONTENT_WIDTH, Inches(0.6))
        tf = desc_box.text_frame
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = "Online/Sosial və ənənəvi media\ndatalarının əsasında təhlil."
        p.font.size = Pt(14.5)
        p.font.color.rgb = RGBColor(0, 123, 191)
        p.line_spacing = Pt(18)
        top += Inches(0.6)

        # Divider Line
        line = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, top, LEFT_CONTENT_WIDTH, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 123, 191)
        line.line.fill.background()
        top += SPACING_MEDIUM

        # "Analitik hesabat"
        report_box = title_slide.shapes.add_textbox(LEFT_MARGIN, top, LEFT_CONTENT_WIDTH, Inches(0.3))
        tf = report_box.text_frame
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = "Analitik hesabat"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0, 123, 191)
        top += Inches(0.3)

        # Date
        date_box = title_slide.shapes.add_textbox(LEFT_MARGIN, top, LEFT_CONTENT_WIDTH, Inches(0.3))
        tf = date_box.text_frame
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = f"{start_date}-{end_date}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 123, 191)
        top += Inches(1.3)

        # --- RIGHT Side Company Logo ---

        if os.path.exists(company_logo_path):
            logo_top = Inches(2.5)
            right_logo_left = RIGHT_HALF_CENTER - LOGO_WIDTH / 2
            title_slide.shapes.add_picture(
                company_logo_path,
                right_logo_left, logo_top,
                width=LOGO_WIDTH
            )

        # endregion

        # region Second slide - Methodology description
        logger.debug("Creating methodology slide")
        method_slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Remove default textbox
        for shape in method_slide.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)

        SLIDE_WIDTH = prs.slide_width
        SLIDE_HEIGHT = prs.slide_height
        HALF_WIDTH = SLIDE_WIDTH / 2

        # Background color
        background = method_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR

        # --- Title on top-left ---
        title_box = method_slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(5), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Metodologiyanın təsviri"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 123, 191)

        # --- NeuroTime logo (top-right corner) ---
        if os.path.exists(neurotime_logo_path):
            method_slide.shapes.add_picture(
                neurotime_logo_path,
                SLIDE_WIDTH - Inches(1.5), Inches(0.3),
                width=Inches(1.2)
            )

        # --- MediaEye logo (centered in left half) ---
        if os.path.exists(mediaeye_logo_path):
            max_logo_width = Inches(4.5)
            max_logo_height = Inches(4.5)

            with Image.open(mediaeye_logo_path) as img:
                width, height = img.size
                aspect_ratio = width / height

                if aspect_ratio > 1:  # Wider than tall
                    display_width = max_logo_width
                    display_height = max_logo_width / aspect_ratio
                else:
                    display_height = max_logo_height
                    display_width = max_logo_height * aspect_ratio

            logo_left = (HALF_WIDTH - display_width) / 2
            logo_top = (SLIDE_HEIGHT - display_height) / 2

            method_slide.shapes.add_picture(
                mediaeye_logo_path,
                logo_left, logo_top,
                width=display_width,
                height=display_height
            )

        # --- Competitor logos grid (up to 25, in right half, centered) ---
        if competitor_logo_paths:
            max_logos = min(25, len(competitor_logo_paths))
            logos_per_row = 5
            logos_per_col = 5

            max_logo_width = Inches(1.1)
            max_logo_height = Inches(0.5)
            spacing = Inches(0.15)

            grid_width = logos_per_row * max_logo_width + (logos_per_row - 1) * spacing
            grid_height = logos_per_col * max_logo_height + (logos_per_col - 1) * spacing

            grid_left_start = HALF_WIDTH + (HALF_WIDTH - grid_width) / 2 - Inches(1)
            grid_top_start = (SLIDE_HEIGHT - grid_height) / 2

            for i, logo_path in enumerate(competitor_logo_paths[:max_logos]):
                if os.path.exists(logo_path):
                    row = i // logos_per_row
                    col = i % logos_per_row
                    left = grid_left_start + col * (max_logo_width + spacing)
                    top = grid_top_start + row * (max_logo_height + spacing)

                    with Image.open(logo_path) as img:
                        width, height = img.size
                        aspect_ratio = width / height

                        if aspect_ratio > 1:
                            display_width = max_logo_width
                            display_height = max_logo_width / aspect_ratio
                        else:
                            display_height = max_logo_height
                            display_width = max_logo_height * aspect_ratio

                    method_slide.shapes.add_picture(
                        logo_path,
                        left + (max_logo_width - display_width) / 2,
                        top + (max_logo_height - display_height) / 2,
                        width=display_width,
                        height=display_height
                    )
        # endregion

        # region Third slide - Grid layout with links and charts
        logger.debug("Creating third slide with grid layout")
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])

        # Remove default textbox
        for shape in slide3.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)

        add_slide_header(slide3, company_logo_path, start_date, end_date, "XƏBƏRLƏRİN ANALİZİ", title_color, template_color)
        add_side_line(slide3, template_color)

        # Set slide background
        background = slide3.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR

        # grid layout positions
        # Top row starts below header
        top_row_y = Inches(1)
        # Bottom row starts lower to accommodate larger top section
        bottom_row_y = Inches(4)
        # Left column starts at left margin
        left_col_x = Inches(0.5)
        # Right column starts at middle of slide
        right_col_x = Inches(6.5)
        # Width for each section
        section_width = Inches(6.5)
        top_section_height = Inches(2.5)
        bottom_section_height = Inches(3)

        # Chart specific widths
        multiline_chart_width = Inches(7)
        donut_chart_width = Inches(5)

        # Adjusted positioning for centered charts
        multiline_chart_x = Inches(0.5)
        donut_chart_x = Inches(8)  # Adjusted to center the smaller donut chart

        # Add positive links section (top left)
        if positive_links:
            left, top = left_col_x, top_row_y
            width = section_width - Inches(0.5)  # Adjusted width to fit better
            height = Inches(0.5)  # Height for individual links
            title_height = Inches(0.3)

            # Add title with split colors
            bg_box = add_bg_box(slide3, left, top, width, top_section_height + Inches(0.2), color=CHART_BG_COLOR)
            title_box = slide3.shapes.add_textbox(left, top, width, title_height)
            tf = title_box.text_frame

            p = tf.add_paragraph()
            # Add "Positiv" in green
            run1 = p.add_run()
            run1.text = "Positiv "
            run1.font.color.rgb = RGBColor(40, 167, 69)  # Green color
            run1.font.size = Pt(16)
            run1.font.bold = True
            # Add "xeberler" in black
            run2 = p.add_run()
            run2.text = "xəbərlərə və saylar aşağıda qeyd olunmuşdur. Xəbərlərə aid nümunələrə başlıqlardan keçid edə bilərsiniz."
            run2.font.color.rgb = RGBColor(0, 0, 0)  # Black color
            run2.font.size = Pt(16)
            run2.font.bold = False
            p.alignment = PP_ALIGN.LEFT
            top += Inches(0.3)  # Spacing after title
            tf.word_wrap = True
            
            for i, link in enumerate(positive_links):
                link_box = slide3.shapes.add_textbox(left, top + height + (height * i), width, height)
                tf = link_box.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = f"🔗 {link}"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 112, 192)  # Blue color
                p.alignment = PP_ALIGN.LEFT
                r = p.runs[0]
                r.hyperlink.address = link


        # Add negative links section (top right)
        if negative_links:
            left, top = right_col_x + Inches(0.5), top_row_y
            width = section_width - Inches(0.5)
            height = Inches(0.5)  # Height for individual links
            title_height = Inches(0.3)
            bg_box = add_bg_box(slide3, left, top, width, top_section_height + Inches(0.2), color=CHART_BG_COLOR)
            # Add title with split colors
            title_box = slide3.shapes.add_textbox(left, top, width, title_height)
            tf = title_box.text_frame
            p = tf.add_paragraph()
            # Add "Negativ" in red
            run1 = p.add_run()
            run1.text = "Negativ "
            run1.font.color.rgb = RGBColor(220, 53, 69)  # Red color
            run1.font.size = Pt(16)
            run1.font.bold = True
            # Add "xeberler" in black
            run2 = p.add_run()
            run2.text = "xəbərlər və saylar aşağıda qeyd olunmuşdur. Xəbərlərə aid nümunələrə başlıqlardan keçid edə bilərsiniz."
            run2.font.color.rgb = RGBColor(0, 0, 0)  # Black color
            run2.font.size = Pt(16)
            run2.font.bold = False
            p.alignment = PP_ALIGN.LEFT
            top += Inches(0.3)  # Spacing after title
            tf.word_wrap = True
            
            for i, link in enumerate(negative_links):
                link_box = slide3.shapes.add_textbox(left, top + height + (height * i), width, height)
                tf = link_box.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = f"🔗 {link}"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(246, 1, 64)  # Red color
                p.alignment = PP_ALIGN.LEFT
                r = p.runs[0]
                r.hyperlink.address = link

        # Get data from combined_sources
        logger.debug("Processing combined sources data")
        if 'combined_sources' not in data_frames:
            raise ValueError("combined_sources Excel file is missing")

        if 'News' not in data_frames['combined_sources']:
            raise ValueError("News sheet is missing in combined_sources Excel file")

        combined_data = data_frames['combined_sources']['News']
        combined_data['Day'] = pd.to_datetime(combined_data['Day'])
        if has_competitors:
            company_data = combined_data[combined_data['Company'] == company_name]
        else:
            company_data = combined_data
        sentiment_data = company_data.groupby(['Day', 'Sentiment']).size().unstack(fill_value=0)
        sentiment_data = sentiment_data.sort_index()
        sentiment_counts = get_sentiment_counts(combined_data)

        # Create multiline chart (bottom left) - wider
        logger.debug("Creating multiline chart")
        chart_data = CategoryChartData()
        chart_data.categories = sentiment_data.index.tolist()

        for sentiment in [1, 0, -1]:
            series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
            if sentiment in sentiment_data.columns:
                chart_data.add_series(series_name, sentiment_data[sentiment].tolist())

        left, top = multiline_chart_x, bottom_row_y
        x, y, cx, cy = left, top, multiline_chart_width, bottom_section_height

        create_sentiment_line_chart(slide3, x, y, cx, cy, chart_data, title="Xəbərlərin sentiment və zamana görə bölgüsü", graph_color=graph_color)

        # Create donut chart (bottom right) - narrower and centered
        logger.debug("Creating donut chart")
        left, top = donut_chart_x, bottom_row_y
        x, y, cx, cy = left, top, donut_chart_width, bottom_section_height
        create_sentiment_donut_chart(slide3, x, y, cx, cy, sentiment_counts, graph_color=graph_color)
        # endregion

        # region Fourth slide - Vertical multibar chart
        if has_competitors:
            logger.debug("Creating Fourth slide with vertical multibar chart")
            slide4 = prs.slides.add_slide(prs.slide_layouts[5])
            # Remove default textbox
            for shape in slide4.shapes:
                if shape.has_text_frame:
                    sp = shape._element
                    sp.getparent().remove(sp)

            add_slide_header(slide4, company_logo_path, start_date, end_date, "XƏBƏRLƏRİN ANALİZİ", title_color, template_color)
            add_side_line(slide4, template_color)

            # Set slide background
            background = slide4.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = SLIDE_BG_COLOR
            
            # Get company sentiment data and sort it by total values
            company_sentiments = get_company_sentiment_counts(combined_data)
            # Calculate total posts for each company
            totals = company_sentiments.sum(axis=1)
            # Sort the sentiment data based on totals
            company_sentiments = company_sentiments.loc[totals.sort_values(ascending=False).index]
            
            chart_data = CategoryChartData()
            chart_data.categories = company_sentiments.index.tolist()
            
            for sentiment in [1, 0, -1]:
                series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                if sentiment in company_sentiments.columns:
                    chart_data.add_series(series_name, company_sentiments[sentiment].tolist())
            
            # Calculate centered position
            # Slide width: 13.33", Chart width: 11"
            # Center horizontally: (13.33 - 11) / 2 = 1.165"
            # Slide height: 7.5", Header height: 0.8", Chart height: 5"
            # Center vertically in remaining space: 0.8 + (6.7 - 5) / 2 = 1.65"
            x = Inches(1.165)
            y = Inches(1.65)
            cx = Inches(11)  # Chart width
            cy = Inches(5)   # Chart height

            bg_box = add_bg_box(slide4, x, y, cx, cy, color=CHART_BG_COLOR)
            chart = slide4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
            
            # Configure legend
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = True
            chart.legend.font.size = Pt(12)
            chart.legend.font.bold = False
            
            # Set chart title with icon
            chart.has_title = True
            chart.chart_title.text_frame.text = "📊 Post saylarına görə bankların bölgüsü"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = False
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = graph_color
            # Set white background for chart
            chart.chart_style = 2  # White background style
            
            # Apply fill to plot area if available and remove gridlines
            try:
                if hasattr(chart, 'plot_area'):
                    chart.plot_area.format.fill.solid()
                    chart.plot_area.format.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    
                # Remove all gridlines
                # Remove major and minor gridlines from the plot
                chart.plots[0].has_major_gridlines = False
                chart.plots[0].has_minor_gridlines = False
                
                # Also remove gridlines from both axes
                if hasattr(chart, 'category_axis') and chart.category_axis:
                    chart.category_axis.has_major_gridlines = False
                    chart.category_axis.has_minor_gridlines = False
                
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    chart.value_axis.has_major_gridlines = False
                    chart.value_axis.has_minor_gridlines = False
                
                # Remove axis lines
                chart.category_axis.format.line.fill.background()
                chart.value_axis.format.line.fill.background()
            except Exception as e:
                logger.warning(f"Could not set plot area formatting: {e}")
            
            # Set diagonal labels for x-axis
            if hasattr(chart, 'category_axis'):
                chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
                chart.category_axis.text_rotation = 90
            
            # Set axis titles and labels
            if hasattr(chart, 'category_axis'):
                chart.category_axis.tick_labels.font.size = Pt(7)
                if hasattr(chart.category_axis, 'axis_title') and chart.category_axis.axis_title:
                    chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)
            if hasattr(chart, 'value_axis'):
                chart.value_axis.tick_labels.font.size = Pt(7)
                if hasattr(chart.value_axis, 'axis_title') and chart.value_axis.axis_title:
                    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)
            
            # Set colors for bar chart
            for i, series in enumerate(chart.series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = SENTIMENT_COLORS[list(SENTIMENT_COLORS.keys())[i]]
                # Add data labels to outside end
                series.has_data_labels = True
                data_labels = series.data_labels
                data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                data_labels.font.size = Pt(10)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                data_labels.font.bold = True
                data_labels.number_format = '0'  # Show whole numbers only
        # endregion

        # region Fifth slide - Author count horizontal bar chart
        logger.debug("Creating fifth slide with author count chart")
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        # Remove default textbox
        for shape in slide5.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)

        add_slide_header(slide5, company_logo_path, start_date, end_date, "Xəbər saylarının saytlar üzərindən bölgüsü", title_color, template_color)
        add_side_line(slide5, template_color)

        # Set slide background
        background = slide5.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR

        # Filter and group data by author
        if has_competitors:
            # Filter by company, then group by Author
            author_data = combined_data[combined_data['Company'] == company_name].groupby('Author').size().sort_values(ascending=True)
        else:
            # Group by Author for all companies
            author_data = combined_data.groupby('Author').size().sort_values(ascending=True)
        author_data = author_data.tail(20)
        # Ensure we have data to prevent errors
        if len(author_data) == 0:
            logger.warning("No author data found for chart")
            # Add a text placeholder or skip chart creation
        else:
            chart_data = CategoryChartData()
            chart_data.categories = author_data.index.tolist()
            chart_data.add_series('', author_data.values.tolist())  # Empty series name to remove title
            
            # Calculate centered position
            # Slide dimensions: 13.33" x 7.5"
            # Header height: 0.8"
            # Available space: 7.5 - 0.8 = 6.7"
            chart_width = Inches(11.5)
            chart_height = Inches(6.0)
            
            # Center horizontally: (13.33 - 11.5) / 2 = 0.915
            # Center vertically in remaining space: 0.8 + (6.7 - 6.0) / 2 = 1.15
            x = Inches(0.9)
            y = Inches(1.15)
            
            bg_box = add_bg_box(slide5, x, y, chart_width, chart_height, color=CHART_BG_COLOR)
            chart = slide5.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, chart_width, chart_height, chart_data).chart
            
            # Remove legend
            chart.has_legend = False
            chart.has_title = False
            
            # Set white background for chart area and plot area
            try:
                # Chart area background
                chart.chart_area.fill.solid()
                chart.chart_area.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                # Plot area background - solid white fill
                try:
                    chart.plot_area.format.fill.solid()
                    chart.plot_area.format.fill.fore_color.rgb = RGBColor(255, 255, 255)
                except Exception as e:
                    logger.warning(f"Could not set plot area fill color: {e}")

            except Exception as e:
                logger.warning(f"Could not set chart background: {e}")
            
            # Remove all gridlines
            try:
                # Remove major and minor gridlines from the plot
                chart.plots[0].has_major_gridlines = False
                chart.plots[0].has_minor_gridlines = False
                
                # Also remove gridlines from both axes
                if hasattr(chart, 'category_axis') and chart.category_axis:
                    chart.category_axis.has_major_gridlines = False
                    chart.category_axis.has_minor_gridlines = False
                
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    chart.value_axis.has_major_gridlines = False
                    chart.value_axis.has_minor_gridlines = False
                    
            except Exception as e:
                logger.warning(f"Could not remove gridlines: {e}")
            
            # Configure category axis (Y-axis for horizontal bar chart) - keep only this axis
            try:
                if hasattr(chart, 'category_axis') and chart.category_axis:
                    chart.category_axis.tick_labels.font.size = Pt(9)
                    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW  # Keep Y-axis labels on left
                    # Remove axis title
                    chart.category_axis.has_title = False
                    # Keep the axis line visible
                    chart.category_axis.format.line.color.rgb = RGBColor(0, 0, 0)  # Black axis line
            except Exception as e:
                logger.warning(f"Could not configure category axis: {e}")
            
            # Remove/hide value axis (X-axis for horizontal bar chart)
            try:
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    # Hide the axis completely
                    chart.value_axis.visible = False
                    # Alternative: make axis line transparent
                    chart.value_axis.format.line.fill.background()
                    # Remove axis title
                    chart.value_axis.has_title = False
                    # Hide tick labels
                    chart.value_axis.tick_labels.font.size = Pt(1)
                    chart.value_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)  # Make invisible
            except Exception as e:
                logger.warning(f"Could not configure value axis: {e}")
            
            # Format bars and add data labels
            try:
                for series in chart.series:
                    # Set bar color
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(160, 208, 255)  # Blue color
                    
                    # Make bars thick but add space between them
                    series.format.gap_width = 40  # Thick bars with some spacing
                    
                    # Remove line around bars to prevent errors
                    series.format.line.fill.background()
                    
                    # Add data labels to outside end
                    series.has_data_labels = True
                    data_labels = series.data_labels
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    data_labels.font.size = Pt(10)
                    data_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                    data_labels.font.bold = True
                    data_labels.number_format = '0'  # Show whole numbers only
                    
            except Exception as e:
                logger.warning(f"Could not format series or data labels: {e}")
            
            # Additional formatting for bars spacing
            try:
                # Access the plot area and modify bar thickness
                plot = chart.plots[0]
                plot.gap_width = 40  # Set plot-level gap width for spacing
            except Exception as e:
                logger.warning(f"Could not set plot formatting: {e}")
            
            # Remove chart title if it exists
            try:
                if chart.has_title:
                    chart.chart_title.text_frame.text = ''
                    chart.chart_title.include_in_layout = False
            except Exception as e:
                logger.warning(f"Could not remove chart title: {e}")

        # endregion

        # region Sixth slide - Facebook metrics and sentiment analysis
        logger.debug("Creating Sixth slide with Facebook metrics")
        slide6 = prs.slides.add_slide(prs.slide_layouts[5])
        # Remove default textbox
        for shape in slide6.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)
        add_slide_header(slide6, company_logo_path, start_date, end_date, "Facebook postlarının analizi", title_color, template_color)
        add_side_line(slide6, template_color)

        # Set slide background
        background = slide6.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR
        
        # Left side - Facebook metrics
        if 'official_facebook' in data_frames and len(data_frames['official_facebook']) > 0:
            fb_data = list(data_frames['official_facebook'].values())[0]
            fb_metrics = fb_data[fb_data['author_name'] == company_name]

        metrics = {
            'Post sayı': len(fb_metrics),
            'Şərh sayı': fb_metrics['comment_count'].sum(),
            'Bəyənmə sayı': fb_metrics['like_count'].sum(),
            'Paylaşım sayı': fb_metrics['share_count'].sum(),
            'Baxış sayı': fb_metrics['view_count'].sum()
        }

        if has_competitors:
            left = Inches(1)
            header_height = Inches(0.8)
            top = header_height + Inches(0.2)  # Start below header with small margin
            metrics_width = Inches(2.2)
            metrics_height = Inches(1.0)
            slide_height = Inches(7.5)  # Standard slide height
            available_height = slide_height - header_height - Inches(0.4)  # Bottom margin
            total_items = len(metrics)

            card_spacing = Inches(0.3)  # Space between cards
            total_cards_height = metrics_height * total_items
            total_spacing_height = card_spacing * (total_items - 1)
            remaining_space = available_height - total_cards_height - total_spacing_height
            top_margin = remaining_space / 2

            for i, (metric, value) in enumerate(metrics.items()):
                box_top = top + top_margin + (metrics_height * i) + (card_spacing * i)
                
                # White background card
                bg_box = add_bg_box(slide6, left, box_top, metrics_width, metrics_height, color=CHART_BG_COLOR)

                # Red icon background (adjusted for new card height)
                icon_width = Inches(0.5)
                icon_height = Inches(0.5)
                icon_left = left + Inches(0.1)
                icon_top = box_top + (metrics_height - icon_height) / 2  # Centered vertically
                
                icon_box = slide6.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    icon_left, icon_top, icon_width, icon_height
                )
                icon_box.fill.solid()
                icon_box.fill.fore_color.rgb = template_color
                icon_box.line.fill.background()
                
                icon_text = slide6.shapes.add_textbox(
                    icon_left, icon_top, icon_width, icon_height
                )
                icon_frame = icon_text.text_frame
                icon_frame.clear()
                icon_frame.margin_top = Inches(0)
                icon_frame.margin_bottom = Inches(0)
                icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
                
                icon_p = icon_frame.paragraphs[0]
                icon_p.text = METRIC_ICONS.get(metric, "📊")  # Use specific icon or default
                icon_p.alignment = PP_ALIGN.CENTER
                icon_p.font.size = Pt(22)  # Adjusted for smaller icon box
                icon_p.font.color.rgb = RGBColor(255, 255, 255)
                
                # Combined value + metric label, vertically centered inside white bg
                text_left = icon_left + icon_width + Inches(0.15)
                text_width = metrics_width - icon_width - Inches(0.25)
                
                value_textbox = slide6.shapes.add_textbox(
                    text_left, box_top, text_width, metrics_height
                )
                tf = value_textbox.text_frame
                tf.clear()
                tf.margin_top = Inches(0)
                tf.margin_bottom = Inches(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center content vertically
                
                # Big number
                value_p = tf.paragraphs[0]
                value_p.text = format_number_with_k(value)
                value_p.alignment = PP_ALIGN.LEFT
                value_p.font.size = Pt(24)
                value_p.font.bold = True
                value_p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray for better readability
                
                # Metric name
                metric_p = tf.add_paragraph()
                metric_p.text = metric
                metric_p.alignment = PP_ALIGN.LEFT
                metric_p.font.size = Pt(11)
                metric_p.font.color.rgb = RGBColor(102, 102, 102)  # Medium gray

            # Right side - Sentiment analysis
            if 'Facebook' in data_frames['combined_sources']:
                fb_sentiment_data = data_frames['combined_sources']['Facebook']
                company_sentiment = fb_sentiment_data[fb_sentiment_data['Company'] == company_name]
                
                # Donut chart
                sentiment_counts = company_sentiment['Sentiment'].value_counts()
                donut_data = ChartData()
                donut_data.categories = ['Positive', 'Neutral', 'Negative']
                donut_data.add_series('', [
                    sentiment_counts.get(1, 0),
                    sentiment_counts.get(0, 0),
                    sentiment_counts.get(-1, 0)
                ])
                
                # Start right section (80% width) layout
                right_section_left = Inches(3.7)  # After left 20% section
                right_width = Inches(9.13)  # ~80% of slide width
                
                # Position donut chart in upper left of right section
                donut_size = Inches(3.5)
                x = right_section_left
                y = Inches(1.2)
                
                create_sentiment_donut_chart(slide6, x, y, donut_size, donut_size - Inches(0.4), sentiment_counts, graph_color=graph_color)

                # Multiline chart
                sentiment_by_date = company_sentiment.groupby('Day')['Sentiment'].value_counts().unstack(fill_value=0)
                chart_data = CategoryChartData()
                chart_data.categories = sentiment_by_date.index.tolist()
                
                for sentiment in [1, 0, -1]:
                    series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                    if sentiment in sentiment_by_date.columns:
                        chart_data.add_series(series_name, sentiment_by_date[sentiment].tolist())
                
                # Position multiline chart to right of donut
                x = right_section_left + donut_size + Inches(0.3)  # After donut chart
                y = Inches(1.2)
                cx = right_width - donut_size - Inches(0.3)  # Remaining width
                cy = donut_size - Inches(0.4)  # Same height as donut
                
                title = "Postların zamana və sentimentə görə bölgüsü"
                create_sentiment_line_chart(slide6, x, y, cx, cy, chart_data, title, graph_color=graph_color)
                    
                # Vertical multibar chart
                x = right_section_left
                y = Inches(4.5)  # Below upper charts
                cx = right_width  # Full width of right section
                cy = Inches(2.5)  # Remaining height
                bg_box = add_bg_box(slide6, x, y, cx, cy, color=CHART_BG_COLOR)

                company_sentiments = fb_sentiment_data.groupby('Company')['Sentiment'].value_counts().unstack(fill_value=0)
                # Sort by total sentiment values
                totals = company_sentiments.sum(axis=1)
                company_sentiments = company_sentiments.loc[totals.sort_values(ascending=False).index]

                chart_data = CategoryChartData()
                chart_data.categories = company_sentiments.index.tolist()
                
                for sentiment in [1, 0, -1]:
                    series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                    if sentiment in company_sentiments.columns:
                        chart_data.add_series(series_name, company_sentiments[sentiment].tolist())
                
                # Position multibar chart to span full width of right section
                chart = slide6.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP
                chart.legend.font.size = Pt(12)

                # Set chart background and formatting
                apply_chart_formatting(chart, title="Post saylarına görə bankların bölgüsü", graph_color=graph_color)
                apply_sentiment_colors(chart)
        else: # no competitiors
            left = Inches(0.5)

            header_height = Inches(0.8)
            metrics_width = Inches(1.8)  # Reduced from 2.2 to 1.5
            metrics_height = Inches(0.85)  # Reduced from 1.0
            total_items = len(metrics)

            # Calculate updated layout with new height
            card_spacing = Inches(0.3)
            total_cards_height = metrics_height * total_items
            total_spacing_height = card_spacing * (total_items - 1)
            card_section_height = total_cards_height + total_spacing_height

            # Maintain cards at bottom by anchoring them to bottom of available area
            top = header_height + Inches(0.2)
            slide_height = Inches(7.5)
            available_height = slide_height - header_height - Inches(0.4)
            top_margin = available_height - card_section_height  # Pushed down to keep bottom position

            # Add top description card
            text_card_height = Inches(0.6)
            text_card_top = top + top_margin - text_card_height - Inches(0.15)  # A bit above the metrics section
            text_card = add_bg_box(slide6, left, text_card_top, metrics_width, text_card_height, color=CHART_BG_COLOR)

            text_box = slide6.shapes.add_textbox(left, text_card_top, metrics_width, text_card_height)
            text_frame = text_box.text_frame
            text_frame.clear()
            text_frame.margin_top = Inches(0)
            text_frame.margin_bottom = Inches(0)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            para = text_frame.paragraphs[0]
            para.text = "📌 Rəsmi səhifənin analizi"
            para.word_wrap = True
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(13)
            para.font.bold = True
            para.font.color.rgb = graph_color

            # Now render each metric card (at bottom)
            for i, (metric, value) in enumerate(metrics.items()):
                box_top = top + top_margin + (metrics_height * i) + (card_spacing * i)

                # --- white bg box ---
                bg_box = add_bg_box(slide6, left, box_top, metrics_width, metrics_height, color=CHART_BG_COLOR)

                # --- red icon box ---
                icon_width = Inches(0.5)
                icon_height = Inches(0.5)
                icon_left = left + Inches(0.1)
                icon_top = box_top + (metrics_height - icon_height) / 2

                icon_box = slide6.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, icon_left, icon_top, icon_width, icon_height
                )
                icon_box.fill.solid()
                icon_box.fill.fore_color.rgb = template_color
                icon_box.line.fill.background()

                icon_text = slide6.shapes.add_textbox(icon_left, icon_top, icon_width, icon_height)
                icon_frame = icon_text.text_frame
                icon_frame.clear()
                icon_frame.margin_top = Inches(0)
                icon_frame.margin_bottom = Inches(0)
                icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                icon_p = icon_frame.paragraphs[0]
                icon_p.text = METRIC_ICONS.get(metric, "📊")
                icon_p.alignment = PP_ALIGN.CENTER
                icon_p.font.size = Pt(22)
                icon_p.font.color.rgb = RGBColor(255, 255, 255)

                # --- value + metric text ---
                text_left = icon_left + icon_width + Inches(0.15)
                text_width = metrics_width - icon_width - Inches(0.25)

                value_textbox = slide6.shapes.add_textbox(text_left, box_top, text_width, metrics_height)
                tf = value_textbox.text_frame
                tf.clear()
                tf.margin_top = Inches(0)
                tf.margin_bottom = Inches(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                value_p = tf.paragraphs[0]
                value_p.text = format_number_with_k(value)
                value_p.alignment = PP_ALIGN.LEFT
                value_p.font.size = Pt(24)
                value_p.font.bold = True
                value_p.font.color.rgb = RGBColor(51, 51, 51)

                metric_p = tf.add_paragraph()
                metric_p.text = metric
                metric_p.alignment = PP_ALIGN.LEFT
                metric_p.font.size = Pt(11)
                metric_p.font.color.rgb = RGBColor(102, 102, 102)
                
            # Middle side - Sentiment analysis
            if 'Facebook' in data_frames['combined_sources']:
                fb_sentiment_data = data_frames['combined_sources']['Facebook']
                company_sentiment = fb_sentiment_data
                
                # Donut chart
                sentiment_counts = company_sentiment['Sentiment'].value_counts()
                donut_data = ChartData()
                donut_data.categories = ['Positive', 'Neutral', 'Negative']
                donut_data.add_series('', [
                    sentiment_counts.get(1, 0),
                    sentiment_counts.get(0, 0),
                    sentiment_counts.get(-1, 0)
                ])
                
                # Start right section with equal margins
                right_section_left = metrics_width + Inches(0.8)  # After left metrics section
                right_width = Inches(8.33)
                
                # Position donut chart in upper left of right section
                donut_size = Inches(3.5)
                x = right_section_left
                y = Inches(1.2)
                
                create_sentiment_donut_chart(slide6, x, y, donut_size, donut_size - Inches(0.4), sentiment_counts, graph_color=graph_color)

                # Text section instead of multiline chart
                x_text = right_section_left + donut_size + Inches(0.3)  # After donut chart
                y_text = Inches(1.2)
                cx_text = right_width - donut_size - Inches(0.3)  # Remaining width
                cy_text = donut_size - Inches(0.4)  # Same height as donut
                bg_box = add_bg_box(slide6, x_text, y_text, cx_text, cy_text, color=CHART_BG_COLOR)

                # --- TITLE TEXTBOX ---
                title_left = x_text + Inches(0.2)
                title_top = y_text + Inches(0.2)
                title_width = cx_text - Inches(0.4)
                title_height = Inches(0.6)
                title_textbox = slide6.shapes.add_textbox(title_left, title_top, title_width, title_height)
                title_tf = title_textbox.text_frame
                title_tf.clear()
                title_tf.word_wrap = True

                title_p = title_tf.paragraphs[0]
                title_p.clear()

                run_icon = title_p.add_run()
                run_icon.text = f"{CHARTS_ICONS['Time Distribution']} Facebook postlarının "
                run_icon.font.size = Pt(16)
                run_icon.font.bold = True
                run_icon.font.color.rgb = RGBColor(0, 0, 0)

                run_positive = title_p.add_run()
                run_positive.text = "pozitiv "
                run_positive.font.size = Pt(16)
                run_positive.font.bold = True
                run_positive.font.color.rgb = RGBColor(40, 167, 69)  # Green

                run_negative = title_p.add_run()
                run_negative.text = "və "
                run_negative.font.size = Pt(16)
                run_negative.font.bold = True
                run_negative.font.color.rgb = RGBColor(0, 0, 0)

                run_negative = title_p.add_run()
                run_negative.text = "neqativ "
                run_negative.font.size = Pt(16)
                run_negative.font.bold = True
                run_negative.font.color.rgb = RGBColor(220, 53, 69)  # Red

                run_text = title_p.add_run()
                run_text.text = "xəbər məzmunları aşağıda qeyd edilmişdir."
                run_text.font.size = Pt(16)
                run_text.font.bold = True
                run_text.font.color.rgb = RGBColor(0, 0, 0)

                title_p.alignment = PP_ALIGN.LEFT

                # --- PARAGRAPH 1 TEXTBOX ---
                para1_left = x_text + Inches(0.2)
                para1_top = title_top + title_height + Inches(0.1)
                para1_width = cx_text - Inches(0.4)
                para1_height = Inches(1)
                para1_textbox = slide6.shapes.add_textbox(para1_left, para1_top, para1_width, para1_height)
                para1_tf = para1_textbox.text_frame
                para1_tf.clear()
                para1_tf.word_wrap = True

                para1 = para1_tf.paragraphs[0]
                para1.level = 0
                para1.clear()

                run1_1 = para1.add_run()
                run1_1.text = 'Prezident İlham Əliyev: “Azəri-Çıraq-Günəşli”, “Abşeron” və “Şahdəniz” bir çox ölkələrin enerji təhlükəsizliyinə töhfə verir; "BP"nin qrantları ilə enerji laboratoriyası qurulacaq; Neft-qaz sektoru üzrə büdcə daxilolmaları 5 %-yə yaxın artıb və bu kimi xəbərlər '
                run1_1.font.size = Pt(12)
                run1_1.font.bold = False
                run1_1.font.color.rgb = RGBColor(51, 51, 51)

                run1_2 = para1.add_run()
                run1_2.text = "pozitiv"
                run1_2.font.size = Pt(12)
                run1_2.font.bold = True
                run1_2.font.color.rgb = RGBColor(40, 167, 69)

                run1_3 = para1.add_run()
                run1_3.text = " olaraq qeyd edilmişdir."
                run1_3.font.size = Pt(12)
                run1_3.font.bold = False
                run1_3.font.color.rgb = RGBColor(51, 51, 51)

                para1.alignment = PP_ALIGN.LEFT

                # --- PARAGRAPH 2 TEXTBOX ---
                para2_left = x_text + Inches(0.2)
                para2_top = para1_top + para1_height + Inches(0.1)
                para2_width = cx_text - Inches(0.4)
                para2_height = Inches(1)
                para2_textbox = slide6.shapes.add_textbox(para2_left, para2_top, para2_width, para2_height)
                para2_tf = para2_textbox.text_frame
                para2_tf.clear()
                para2_tf.word_wrap = True

                para2 = para2_tf.paragraphs[0]
                para2.level = 0
                para2.clear()

                run2_1 = para2.add_run()
                run2_1.text = "Neftin ucuzlaşması səbəbilə ilk 3 ayda şirkətin mənfəəti 49% aşağı düşüb; Can dedik, kardeş dedik Azerbaycan'a şimdi ne oluyor?; Aslında SOCAR ile işbirliği yapan BP'yi sayarsak Tek millet 4 devlet olmuşlar və bu kimi xəbərlər "
                run2_1.font.size = Pt(12)
                run2_1.font.bold = False
                run2_1.font.color.rgb = RGBColor(51, 51, 51)

                run2_2 = para2.add_run()
                run2_2.text = "neqativ"
                run2_2.font.size = Pt(12)
                run2_2.font.bold = True
                run2_2.font.color.rgb = RGBColor(220, 53, 69)

                run2_3 = para2.add_run()
                run2_3.text = " olaraq qeyd edilmişdir."
                run2_3.font.size = Pt(12)
                run2_3.font.bold = False
                run2_3.font.color.rgb = RGBColor(51, 51, 51)

                para2.alignment = PP_ALIGN.LEFT

                # Vertical multibar chart grouped by Day
                x = right_section_left
                y = Inches(4.5)  # Below upper charts
                cx = right_width  # Full width of right section
                cy = Inches(2.5)  # Remaining height
                bg_box = add_bg_box(slide6, x, y, cx, cy, color=CHART_BG_COLOR)

                # Group data by Day instead of Company
                sentiment_by_day = fb_sentiment_data.groupby('Day')['Sentiment'].value_counts().unstack(fill_value=0)
                # Sort by date
                sentiment_by_day = sentiment_by_day.sort_index()

                chart_data = CategoryChartData()
                chart_data.categories = sentiment_by_day.index.tolist()
                
                for sentiment in [1, 0, -1]:
                    series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                    if sentiment in sentiment_by_day.columns:
                        chart_data.add_series(series_name, sentiment_by_day[sentiment].tolist())

                chart = slide6.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP
                chart.legend.font.size = Pt(12)
                chart.has_title = True
                chart.chart_title.text_frame.text = f"{CHARTS_ICONS['Sentiment Trend']} Postların günlər üzrə bölgüsü"
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
                chart.chart_title.text_frame.paragraphs[0].font.bold = False
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = graph_color

                # Set chart background and formatting
                apply_chart_formatting(chart, title="Postların günlər üzrə bölgüsü", icon=CHARTS_ICONS['Sentiment Trend'], graph_color=graph_color)
                for i, series in enumerate(chart.series):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = SENTIMENT_COLORS[list(SENTIMENT_COLORS.keys())[i]]
                    series.has_data_labels = True
                    series.data_labels.font.size = Pt(10)
                    series.data_labels.font.bold = True
                    series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

            # right side - Metrics
            if 'facebook_reachs' in data_frames and len(data_frames['facebook_reachs']) > 0:
                fb_data = list(data_frames['facebook_reachs'].values())[0]
                fb_metrics = fb_data

            metrics = {
                'Post sayı': len(fb_metrics),
                'Şərh sayı': fb_metrics['comment_count'].sum(),
                'Bəyənmə sayı': fb_metrics['like_count'].sum(),
                'Paylaşım sayı': fb_metrics['share_count'].sum(),
                'Baxış sayı': fb_metrics['view_count'].sum()
            }

            # Add top description card
            text_card_height = Inches(0.6)
            right_metrics_left = left + metrics_width + right_width + Inches(0.6)
            text_card_top = top + top_margin - text_card_height - Inches(0.15)  # A bit above the metrics section
            text_card = add_bg_box(slide6, right_metrics_left, text_card_top, metrics_width, text_card_height, color=CHART_BG_COLOR)

            text_box = slide6.shapes.add_textbox(right_metrics_left, text_card_top, metrics_width, text_card_height)
            text_frame = text_box.text_frame
            text_frame.clear()
            text_frame.margin_top = Inches(0)
            text_frame.margin_bottom = Inches(0)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            para = text_frame.paragraphs[0]
            para.text = "📌 Açar sözlü postlar"
            para.word_wrap = True
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(13)
            para.font.bold = True
            para.font.color.rgb = graph_color

            for i, (metric, value) in enumerate(metrics.items()):
                box_top = top + top_margin + (metrics_height * i) + (card_spacing * i)
                # White background card
                bg_box = add_bg_box(slide6, right_metrics_left, box_top, metrics_width, metrics_height, color=CHART_BG_COLOR)

                # Red icon background (adjusted for new card height)
                icon_width = Inches(0.5)
                icon_height = Inches(0.5)
                icon_left = right_metrics_left + Inches(0.1)
                icon_top = box_top + (metrics_height - icon_height) / 2  # Centered vertically
                
                icon_box = slide6.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    icon_left, icon_top, icon_width, icon_height
                )
                icon_box.fill.solid()
                icon_box.fill.fore_color.rgb = template_color
                icon_box.line.fill.background()  # Remove border
                
                # Icon text (white, centered inside red bg)
                icon_text = slide6.shapes.add_textbox(
                    icon_left, icon_top, icon_width, icon_height
                )
                icon_frame = icon_text.text_frame
                icon_frame.clear()
                icon_frame.margin_top = Inches(0)
                icon_frame.margin_bottom = Inches(0)
                icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
                
                icon_p = icon_frame.paragraphs[0]
                icon_p.text = METRIC_ICONS.get(metric, "📊")  # Use specific icon or default
                icon_p.alignment = PP_ALIGN.CENTER
                icon_p.font.size = Pt(22)  # Adjusted for smaller icon box
                icon_p.font.color.rgb = RGBColor(255, 255, 255)
                
                # Combined value + metric label, vertically centered inside white bg
                text_left = icon_left + icon_width + Inches(0.15)
                text_width = metrics_width - icon_width - Inches(0.25)
                
                value_textbox = slide6.shapes.add_textbox(
                    text_left, box_top, text_width, metrics_height
                )
                tf = value_textbox.text_frame
                tf.clear()
                tf.margin_top = Inches(0)
                tf.margin_bottom = Inches(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center content vertically
                
                # Big number
                value_p = tf.paragraphs[0]
                value_p.text = format_number_with_k(value)
                value_p.alignment = PP_ALIGN.LEFT
                value_p.font.size = Pt(24)
                value_p.font.bold = True
                value_p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray for better readability
                
                # Metric name
                metric_p = tf.add_paragraph()
                metric_p.text = metric
                metric_p.alignment = PP_ALIGN.LEFT
                metric_p.font.size = Pt(11)
                metric_p.font.color.rgb = RGBColor(102, 102, 102)  # Medium gray

        # endregion

        # region Seventh slide - Facebook metrics table Combines sources
        if has_competitors:
            logger.debug("Creating eighth slide with Facebook metrics table by company")
            slide7 = prs.slides.add_slide(prs.slide_layouts[5])
            # Remove default textbox
            for shape in slide7.shapes:
                if shape.has_text_frame:
                    sp = shape._element
                    sp.getparent().remove(sp)
            add_slide_header(slide7, company_logo_path, start_date, end_date, "Banklar haqqında paylaşılan postların analizi", title_color, template_color)
            add_side_line(slide7, template_color)

            # Set slide background
            background = slide7.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = SLIDE_BG_COLOR
            
            if 'combined_sources' in data_frames and 'Facebook' in data_frames['combined_sources']:
                fb_data = data_frames['combined_sources']['Facebook']
                
                # Group data by Company
                grouped_data = fb_data.groupby('Company').agg({
                    'comment_count': 'sum',
                    'like_count': 'sum',
                    'share_count': 'sum',
                    'view_count': 'sum'
                }).reset_index()
                
                # Add post count
                post_counts = fb_data.groupby('Company').size().reset_index(name='post_count')
                grouped_data = grouped_data.merge(post_counts, on='Company')
                
                # Create table
                rows = len(grouped_data) + 1  # +1 for header
                cols = 6  # Company, post_count, comment_count, like_count, share_count, view_count
                
                # Calculate centered position
                # Slide width: 13.33", Table width: 12.33" (leaving 0.5" on each side)
                # Center horizontally: (13.33 - 12.33) / 2 = 0.5"
                left = Inches(0.5)
                
                # Slide height: 7.5", Header: 0.8", Table height: 5"
                # Center vertically in remaining space: 0.8 + (6.7 - 5) / 2 = 1.65"
                top = Inches(1)
                width = Inches(12.33)
                height = Inches(5)
                
                table = slide7.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set column widths with wider spacing
                table.columns[0].width = Inches(3)  # Company name (wider for longer names)
                for i in range(1, cols):
                    table.columns[i].width = Inches(1.866)  # Metrics (remaining width distributed evenly)
                
                # Add headers with red background
                headers = ['Banklar ', 'Post sayı', 'Şərh sayı', 'Bəyənmə sayı', 'Paylaşım sayı', 'Baxış sayı']
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = template_color
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
                
                # Add data with alternating row colors
                for i, row in grouped_data.iterrows():
                    for j in range(cols):
                        cell = table.cell(i + 1, j)
                        if j == 0:
                            cell.text = str(row['Company'])
                        elif j == 1:
                            cell.text = str(row['post_count'])
                        elif j == 2:
                            cell.text = str(row['view_count'])
                        elif j == 3:
                            cell.text = str(row['comment_count'])
                        elif j == 4:
                            cell.text = str(row['like_count'])
                        else:
                            cell.text = str(row['share_count'])
                        
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
                        
                        # Set alternating row colors
                        if i % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        # endregion

        # region Eighth slide - Official Facebook metrics table
        if has_competitors:
            logger.debug("Creating seventh slide with Facebook metrics table")
            slide8 = prs.slides.add_slide(prs.slide_layouts[5])
            # Remove default textbox
            for shape in slide8.shapes:
                if shape.has_text_frame:
                    sp = shape._element
                    sp.getparent().remove(sp)
            add_slide_header(slide8, company_logo_path, start_date, end_date, "Bankların rəsmi Facebook səhifələrinin analizi", title_color, template_color)
            add_side_line(slide8, template_color)

            # Set slide background
            background = slide8.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = SLIDE_BG_COLOR
            
            if 'official_facebook' in data_frames and len(data_frames['official_facebook']) > 0:
                fb_data = list(data_frames['official_facebook'].values())[0]
                
                # Verify required columns exist
                required_columns = ['author_name', 'comment_count', 'like_count', 'share_count', 'view_count']
                if not all(col in fb_data.columns for col in required_columns):
                    logger.error("Missing required columns in Facebook data")
                    raise ValueError("Facebook data is missing required columns")
                
                # Group data by author_name
                grouped_data = fb_data.groupby('author_name').agg({
                    'comment_count': 'sum',
                    'like_count': 'sum',
                    'share_count': 'sum',
                    'view_count': 'sum'
                }).reset_index()
                
                # Add post count
                post_counts = fb_data.groupby('author_name').size().reset_index(name='post_count')
                grouped_data = grouped_data.merge(post_counts, on='author_name')
                
                # Create table
                rows = len(grouped_data) + 1  # +1 for header
                cols = 6  # Company, post_count, comment_count, like_count, share_count, view_count
                
                # Calculate centered position
                # Slide width: 13.33", Table width: 12.33" (leaving 0.5" on each side)
                # Center horizontally: (13.33 - 12.33) / 2 = 0.5"
                left = Inches(0.5)
                
                # Slide height: 7.5", Header: 0.8", Table height: 5"
                # Center vertically in remaining space: 0.8 + (6.7 - 5) / 2 = 1.65"
                top = Inches(1)
                width = Inches(12.33)
                height = Inches(5)
                
                table = slide8.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set column widths with wider spacing
                table.columns[0].width = Inches(3)  # Company name (wider for longer names)
                for i in range(1, cols):
                    table.columns[i].width = Inches(1.866)  # Metrics (remaining width distributed evenly)
                
                # Add headers with red background
                headers = ['Banklar ', 'Post sayı', 'Şərh sayı', 'Bəyənmə sayı', 'Paylaşım sayı', 'Baxış sayı']
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = template_color
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
                
                # Add data with alternating row colors
                for i, row in grouped_data.iterrows():
                    for j in range(cols):
                        cell = table.cell(i + 1, j)
                        if j == 0:
                            cell.text = str(row['author_name'])
                        elif j == 1:
                            cell.text = str(row['post_count'])
                        elif j == 2:
                            cell.text = str(row['comment_count'])  # Changed order to match headers
                        elif j == 3:
                            cell.text = str(row['like_count'])    # Changed order to match headers
                        elif j == 4:
                            cell.text = str(row['share_count'])   # Changed order to match headers
                        else:
                            cell.text = str(row['view_count'])    # Changed order to match headers
                        
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
                        
                        # Set alternating row colors
                        if i % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        # endregion

        # region Ninth slide - Instagram metrics and sentiment analysis
        logger.debug("Creating ninth slide with Instagram metrics and sentiment analysis")
        slide9 = prs.slides.add_slide(prs.slide_layouts[5])
        # Remove default textbox
        for shape in slide9.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)
        add_slide_header(slide9, company_logo_path, start_date, end_date, "Instagram postlarının analizi", title_color, template_color)
        add_side_line(slide9, template_color)

        # Set slide background
        background = slide9.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR
        
        # Left side - Instagram metrics from official_instagram
        if 'official_instagram' in data_frames and len(data_frames['official_instagram']) > 0:
            insta_data = list(data_frames['official_instagram'].values())[0]
            if has_competitors:
                company_metrics = insta_data[insta_data['Company'] == company_name]
            else:
                company_metrics = insta_data
            # Icon mapping for network metrics

        network_metrics = {
            'Bəyənmə sayı': company_metrics['Likes'].sum(),
            'Şərh sayı': company_metrics['Comments'].sum()
        }

        # Set up left section (20% width) for metrics
        left = Inches(0.8)
        header_height = Inches(0.8)
        top = header_height + Inches(0.2)  # Start below header with small margin
        metrics_width = Inches(2.7)  # ~20% of slide width (13.33")
        metrics_height = Inches(1.0)
        slide_height = Inches(7.5)  # Standard slide height
        available_height = slide_height - header_height - Inches(0.4)  # Bottom margin
        total_items = len(network_metrics)

        card_spacing = Inches(0.3)  # Space between cards
        total_cards_height = metrics_height * total_items
        total_spacing_height = card_spacing * (total_items - 1)
        remaining_space = available_height - total_cards_height - total_spacing_height
        top_margin = remaining_space / 2

        for i, (metric, value) in enumerate(network_metrics.items()):
            # Calculate position for current metric box
            box_top = top + top_margin + (metrics_height * i) + (card_spacing * i)
            
            # Add white background box for the whole metric
            bg_box = add_bg_box(slide9, left, box_top, metrics_width, metrics_height, color=CHART_BG_COLOR)

            # Add red background for icon (square, centered vertically)
            icon_width = Inches(0.5)
            icon_height = Inches(0.5)
            icon_left = left + Inches(0.1)
            icon_top = box_top + (metrics_height - icon_height) / 2  # Centered vertically
            
            icon_box = slide9.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                icon_left, icon_top, icon_width, icon_height
            )
            icon_box.fill.solid()
            icon_box.fill.fore_color.rgb = template_color
            icon_box.line.fill.background()  # Remove border
            
            # Add icon text (centered inside red bg)
            icon_text = slide9.shapes.add_textbox(
                icon_left, icon_top, icon_width, icon_height
            )
            icon_frame = icon_text.text_frame
            icon_frame.clear()
            icon_frame.margin_top = Inches(0)
            icon_frame.margin_bottom = Inches(0)
            icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
            
            icon_p = icon_frame.paragraphs[0]
            icon_p.text = METRIC_ICONS.get(metric, "📊")  # Use specific icon or default
            icon_p.alignment = PP_ALIGN.CENTER
            icon_p.font.size = Pt(22)  # Adjusted for icon box size
            icon_p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Combined value + metric label, vertically centered inside white bg
            text_left = icon_left + icon_width + Inches(0.15)
            text_width = metrics_width - icon_width - Inches(0.25)
            
            value_textbox = slide9.shapes.add_textbox(
                text_left, box_top, text_width, metrics_height
            )
            tf = value_textbox.text_frame
            tf.clear()
            tf.margin_top = Inches(0)
            tf.margin_bottom = Inches(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center content vertically
            
            # Big number
            value_p = tf.paragraphs[0]
            value_p.text = format_number_with_k(value)
            value_p.alignment = PP_ALIGN.LEFT
            value_p.font.size = Pt(24)  # Adjusted for smaller cards
            value_p.font.bold = True
            value_p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray for better readability
            
            # Metric name
            metric_p = tf.add_paragraph()
            metric_p.text = metric
            metric_p.alignment = PP_ALIGN.LEFT
            metric_p.font.size = Pt(11)  # Adjusted for smaller cards
            metric_p.font.color.rgb = RGBColor(102, 102, 102)  # Medium gray
            
        # Right side - Sentiment analysis from combined_sources
        if 'combined_sources' in data_frames and 'Instagram' in data_frames['combined_sources']:
            insta_sentiment_data = data_frames['combined_sources']['Instagram']
            if has_competitors:
                company_sentiment = insta_sentiment_data[insta_sentiment_data['Company'] == company_name]
            else:
                company_sentiment = insta_sentiment_data
            
            sentiment_counts = company_sentiment['Sentiment'].value_counts()

            # Right section (80% width) layout
            right_section_left = Inches(3.7)  # After left 20% section
            right_width = Inches(9.13)  # 80% of slide width
            
            if has_competitors:
                # Position smaller donut chart centered in its section
                donut_size = Inches(3.5)
                x = right_section_left
                y = Inches(1.2)

                create_sentiment_donut_chart(slide9, x, y, donut_size + Inches(0.3), donut_size - Inches(0.4), sentiment_counts, graph_color=graph_color)

                # Multiline chart
                x = right_section_left + donut_size + Inches(0.5)  # After donut chart
                y = Inches(1.2)
                cx = right_width - donut_size - Inches(0.5)  # Remaining width
                cy = donut_size - Inches(0.4)  # Same height as donut
                sentiment_by_date = company_sentiment.groupby('Day')['Sentiment'].value_counts().unstack(fill_value=0)
                chart_data = CategoryChartData()
                chart_data.categories = sentiment_by_date.index.tolist()
                for sentiment in [1, 0, -1]:
                    series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                    if sentiment in sentiment_by_date.columns:
                        chart_data.add_series(series_name, sentiment_by_date[sentiment].tolist())
                
                title = "Postların zamana və sentimentə görə bölgüsü"
                create_sentiment_line_chart(slide9, x, y, cx, cy, chart_data, title, icon=CHARTS_ICONS['Sentiment Trend'], graph_color=graph_color)
                    
                # Vertical multibar chart for Instagram company sentiment comparison
                # Position multibar chart to take full width of right section
                x = right_section_left
                y = Inches(4.5)  # Position below the top charts
                cx = right_width  # Full width of right section
                cy = Inches(2.5)  # Height
                bg_box = add_bg_box(slide9, x, y, cx, cy, color=CHART_BG_COLOR)

                insta_sentiment_data = insta_sentiment_data[insta_sentiment_data['Sentiment'].isin([-1, 0, 1])]
                company_sentiments = insta_sentiment_data.groupby('Company')['Sentiment'].value_counts().unstack(fill_value=0)
                # Sort by total sentiment values
                totals = company_sentiments.sum(axis=1)
                company_sentiments = company_sentiments.loc[totals.sort_values(ascending=False).index]

                chart_data = CategoryChartData()
                chart_data.categories = company_sentiments.index.tolist()

                for sentiment in [1, 0, -1]:
                    series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                    if sentiment in company_sentiments.columns:
                        chart_data.add_series(series_name, company_sentiments[sentiment].tolist())

                chart = slide9.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP
                chart.legend.font.size = Pt(12)

                # Set chart background and formatting
                apply_sentiment_colors(chart)
                apply_chart_formatting(chart, title="Post saylarına görə bankların bölgüsü", graph_color=graph_color)

            else:
                # Position smaller donut chart centered in its section
                donut_size = Inches(3.5)
                x = right_section_left + right_width - donut_size - Inches(0.3)  # Centered in right section
                y = Inches(1.2)

                create_sentiment_donut_chart(slide9, x, y, donut_size + Inches(0.3), donut_size - Inches(0.4), sentiment_counts, graph_color=graph_color)
                                                        
                # Info section with separated text boxes
                x = right_section_left
                y = Inches(1.2)
                cx = right_width - donut_size - Inches(0.5)
                cy = donut_size - Inches(0.4)  # Same height as donut
                bg_box = add_bg_box(slide9, x, y, cx, cy, color=CHART_BG_COLOR)

                # Title text box
                title_left = x + Inches(0.2)
                title_top = y + Inches(0.2)
                title_width = cx - Inches(0.4)
                title_height = Inches(1.0)  # Fixed height for title
                title_textbox = slide9.shapes.add_textbox(title_left, title_top, title_width, title_height)
                title_tf = title_textbox.text_frame
                title_tf.clear()
                title_tf.word_wrap = True

                # Title paragraph with colored text
                title_p = title_tf.paragraphs[0]
                title_p.clear()

                # Add icon and styled title
                run_icon = title_p.add_run()
                run_icon.text = f"{CHARTS_ICONS['Time Distribution']} Instagram postlarının "
                run_icon.font.size = Pt(16)
                run_icon.font.bold = True
                run_icon.font.color.rgb = RGBColor(0, 0, 0)

                run_positive = title_p.add_run()
                run_positive.text = "Pozitiv "
                run_positive.font.size = Pt(16)
                run_positive.font.bold = True
                run_positive.font.color.rgb = RGBColor(40, 167, 69)  # Green

                run_text = title_p.add_run()
                run_text.text = "xəbər məzmunları aşağıda qeyd edilmişdir. "
                run_text.font.size = Pt(16)
                run_text.font.bold = True
                run_text.font.color.rgb = RGBColor(0, 0, 0)  # Black

                run_negative = title_p.add_run()
                run_negative.text = "Neqativ "
                run_negative.font.size = Pt(16)
                run_negative.font.bold = True
                run_negative.font.color.rgb = RGBColor(220, 53, 69)  # Red

                run_text2 = title_p.add_run()
                run_text2.text = "xəbərlər isə qeydə alınmamışdır."
                run_text2.font.size = Pt(16)
                run_text2.font.bold = True
                run_text2.font.color.rgb = RGBColor(0, 0, 0)  # Black

                title_p.alignment = PP_ALIGN.LEFT

                # First paragraph text box with bullet
                para1_left = x + Inches(0.2)
                para1_top = title_top + title_height + Inches(0.1)
                para1_width = cx - Inches(0.4)
                para1_height = Inches(0.8)  # Adjusted height for content
                para1_textbox = slide9.shapes.add_textbox(para1_left, para1_top, para1_width, para1_height)
                para1_tf = para1_textbox.text_frame
                para1_tf.clear()
                para1_tf.word_wrap = True

                para1 = para1_tf.paragraphs[0]
                para1.level = 0  # Top-level bullet
                para1.text = (
                    "Prezident İlham Əliyev: “Azəri-Çıraq-Günəşli”, “Abşeron” və “Şahdəniz” bir çox ölkələrin enerji təhlükəsizliyinə töhfə verir və bu kimi xəbər pozitiv olaraq qeyd edilmişdir."                )
                para1.font.size = Pt(12)
                para1.font.bold = False
                para1.font.color.rgb = RGBColor(51, 51, 51)
                para1.alignment = PP_ALIGN.LEFT

                # Second paragraph text box with bullet and colored "neqativ"
                para2_left = x + Inches(0.2)
                para2_top = para1_top + para1_height + Inches(0.1)
                para2_width = cx - Inches(0.4)
                para2_height = Inches(0.8)  # Adjusted height for content
                para2_textbox = slide9.shapes.add_textbox(para2_left, para2_top, para2_width, para2_height)
                para2_tf = para2_textbox.text_frame
                para2_tf.clear()
                para2_tf.word_wrap = True

                para2 = para2_tf.paragraphs[0]
                para2.level = 0  # Top-level bullet
                para2.clear()

                # Add text with colored "neqativ" word
                run1 = para2.add_run()
                run1.text = "Azərbaycanda 2025-ci ilin birinci rübündə 6,9 milyon ton neft (kondensatla) hasil edilib ki, bu da ötən illə müqayisədə 5,48% azdır. 5,7 milyon ton ixrac edilib və bu kimi xəbər "
                run1.font.size = Pt(12)
                run1.font.bold = False
                run1.font.color.rgb = RGBColor(51, 51, 51)

                run_neqativ = para2.add_run()
                run_neqativ.text = "neqativ"
                run_neqativ.font.size = Pt(12)
                run_neqativ.font.bold = True
                run_neqativ.font.color.rgb = RGBColor(220, 53, 69)  # Red

                run2 = para2.add_run()
                run2.text = " olaraq qeyd edilmişdir."
                run2.font.size = Pt(12)
                run2.font.bold = False
                run2.font.color.rgb = RGBColor(51, 51, 51)

                para2.alignment = PP_ALIGN.LEFT

                # Vertical stacked bar chart for Instagram company sentiment comparison
                # Position chart to take full width of right section
                x = right_section_left
                y = Inches(4.5)  # Position below the top charts
                cx = right_width  # Full width of right section
                cy = Inches(2.5)  # Height
                bg_box = add_bg_box(slide9, x, y, cx, cy, color=CHART_BG_COLOR)

                # Filter relevant sentiments
                insta_sentiment_data = insta_sentiment_data[insta_sentiment_data['Sentiment'].isin([-1, 0, 1])]

                # Group and reshape the sentiment data
                company_sentiments = insta_sentiment_data.groupby('Day')['Sentiment'].value_counts().unstack(fill_value=0)

                # Ensure all expected sentiment columns exist, even if some are missing
                for sentiment in [-1, 0, 1]:
                    if sentiment not in company_sentiments.columns:
                        company_sentiments[sentiment] = 0

                # Sort by total sentiment counts
                totals = company_sentiments.sum(axis=1)
                company_sentiments = company_sentiments.loc[totals.sort_values(ascending=False).index]

                # Build the chart data
                chart_data = CategoryChartData()
                chart_data.categories = company_sentiments.index.tolist()

                for sentiment in [1, 0, -1]:
                    series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                    chart_data.add_series(series_name, company_sentiments[sentiment].tolist())

                # Add stacked column chart
                chart = slide9.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_STACKED,
                    x, y, cx, cy, chart_data
                ).chart

                # Configure chart appearance
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP
                chart.legend.font.size = Pt(12)

                # Apply custom formatting and sentiment colors
                apply_sentiment_colors(chart)
                apply_chart_formatting(chart, title="Post saylarına görə bankların bölgüsü", graph_color=graph_color)

        # endregion

        # region Tenth slide - Twitter sentiment analysis
        if not has_competitors:
            logger.debug("Creating tenth slide with Linkedin sentiment analysis")
            slide10 = prs.slides.add_slide(prs.slide_layouts[5])
            # Remove default textbox
            for shape in slide10.shapes:
                if shape.has_text_frame:
                    sp = shape._element
                    sp.getparent().remove(sp)
            add_slide_header(slide10, company_logo_path, start_date, end_date, "Twitter postlarının analizi", title_color, template_color)
            add_side_line(slide10, template_color)

            # Set slide background
            background = slide10.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = SLIDE_BG_COLOR
            
            if 'combined_sources' in data_frames and 'Twitter' in data_frames['combined_sources']:
                linkedin_data = data_frames['combined_sources']['Twitter']
                if has_competitors:
                    company_sentiment = linkedin_data[linkedin_data['Company'] == company_name]
                else:
                    company_sentiment = linkedin_data
                
                if not company_sentiment.empty:
                    # Calculate heights accounting for header
                    available_height = Inches(7.5 - HEADER_HEIGHT)  # Total height minus header
                    half_height = available_height / 2
                    full_content_width = Inches(12.33)

                    # Top half section for text and donut charts
                    donut_width = Inches(5)
                    donut_heigth = Inches(3.1)
                    
                    # Text section on left top half (where donut was)
                    x_text = Inches(0.5)
                    y_text = Inches(1.2)  # Just below header
                    cx_text = full_content_width - donut_width - Inches(0.5)  # Remaining width after donut
                    cy_text = Inches(3.1)  # Height for text section
                    
                    try:
                        # Add background box for text section
                        bg_box = add_bg_box(slide10, x_text, y_text, cx_text, cy_text, color=CHART_BG_COLOR)

                        # --- TITLE TEXTBOX ---
                        title_left = x_text + Inches(0.2)
                        title_top = y_text + Inches(0.2)
                        title_width = cx_text - Inches(0.4)
                        title_height = Inches(0.6)
                        title_textbox = slide10.shapes.add_textbox(title_left, title_top, title_width, title_height)
                        title_tf = title_textbox.text_frame
                        title_tf.clear()
                        title_tf.word_wrap = True

                        title_p = title_tf.paragraphs[0]
                        title_p.clear()

                        run_icon = title_p.add_run()
                        run_icon.text = f"{CHARTS_ICONS['Time Distribution']} Twitter postlarının "
                        run_icon.font.size = Pt(16)
                        run_icon.font.bold = True
                        run_icon.font.color.rgb = RGBColor(0, 0, 0)

                        run_positive = title_p.add_run()
                        run_positive.text = "pozitiv "
                        run_positive.font.size = Pt(16)
                        run_positive.font.bold = True
                        run_positive.font.color.rgb = RGBColor(40, 167, 69)  # Green

                        run_negative = title_p.add_run()
                        run_negative.text = "və neqativ "
                        run_negative.font.size = Pt(16)
                        run_negative.font.bold = True
                        run_negative.font.color.rgb = RGBColor(220, 53, 69)  # Red

                        run_text = title_p.add_run()
                        run_text.text = "xəbər məzmunları aşağıda qeyd edilmişdir."
                        run_text.font.size = Pt(16)
                        run_text.font.bold = True
                        run_text.font.color.rgb = RGBColor(0, 0, 0)

                        title_p.alignment = PP_ALIGN.LEFT

                        # --- PARAGRAPH 1 TEXTBOX ---
                        para1_left = x_text + Inches(0.2)
                        para1_top = title_top + title_height + Inches(0.1)
                        para1_width = cx_text - Inches(0.4)
                        para1_height = Inches(1)
                        para1_textbox = slide10.shapes.add_textbox(para1_left, para1_top, para1_width, para1_height)
                        para1_tf = para1_textbox.text_frame
                        para1_tf.clear()
                        para1_tf.word_wrap = True

                        para1 = para1_tf.paragraphs[0]
                        para1.level = 0
                        para1.clear()

                        run1_1 = para1.add_run()
                        run1_1.text = 'BP şirkəti "Azərbaycan Biznes Keys 2025" yarışmasına sponsorluq edib; Prezident İlham Əliyev: “Azəri-Çıraq-Günəşli”, “Abşeron” və “Şahdəniz” bir çox ölkələrin enerji təhlükəsizliyinə töhfə verir; Türkiyə regional əməkdaşlığa çox böyük əhəmiyyət verir və bu kimi xəbərlər '
                        run1_1.font.size = Pt(12)
                        run1_1.font.bold = False
                        run1_1.font.color.rgb = RGBColor(51, 51, 51)

                        run1_2 = para1.add_run()
                        run1_2.text = "pozitiv"
                        run1_2.font.size = Pt(12)
                        run1_2.font.bold = True
                        run1_2.font.color.rgb = RGBColor(40, 167, 69)

                        run1_3 = para1.add_run()
                        run1_3.text = " olaraq qeyd edilmişdir."
                        run1_3.font.size = Pt(12)
                        run1_3.font.bold = False
                        run1_3.font.color.rgb = RGBColor(51, 51, 51)

                        para1.alignment = PP_ALIGN.LEFT

                        # --- PARAGRAPH 2 TEXTBOX ---
                        para2_left = x_text + Inches(0.2)
                        para2_top = para1_top + para1_height + Inches(0.1)
                        para2_width = cx_text - Inches(0.4)
                        para2_height = Inches(1)
                        para2_textbox = slide10.shapes.add_textbox(para2_left, para2_top, para2_width, para2_height)
                        para2_tf = para2_textbox.text_frame
                        para2_tf.clear()
                        para2_tf.word_wrap = True

                        para2 = para2_tf.paragraphs[0]
                        para2.level = 0
                        para2.clear()

                        run2_1 = para2.add_run()
                        run2_1.text = '2015-ci ildən 2020-ci ilə qədər Azərbaycanın “Şahdəniz” qaz layihəsinin 20%-nə sahib olan “Lukoyl” Aİ-yə qaz nəql edərək, hazırda Ukraynada müharibə aparmaq üçün istifadə olunan 63,8 MİLYARD dollar vergini Rusiyaya ödəyib və bu kimi xəbərlər '
                        run2_1.font.size = Pt(12)
                        run2_1.font.bold = False
                        run2_1.font.color.rgb = RGBColor(51, 51, 51)

                        run2_2 = para2.add_run()
                        run2_2.text = "neqativ"
                        run2_2.font.size = Pt(12)
                        run2_2.font.bold = True
                        run2_2.font.color.rgb = RGBColor(220, 53, 69)

                        run2_3 = para2.add_run()
                        run2_3.text = " olaraq qeyd edilmişdir."
                        run2_3.font.size = Pt(12)
                        run2_3.font.bold = False
                        run2_3.font.color.rgb = RGBColor(51, 51, 51)

                        para2.alignment = PP_ALIGN.LEFT

                    except Exception as e:
                        print(f"Error creating text section: {e}")
                    
                    # Donut chart on right top half (where multiline was)
                    x_donut = full_content_width - donut_width + Inches(0.5) # Start after text section
                    y_donut = Inches(1.2)  # Same vertical alignment as text
                    create_sentiment_donut_chart(slide10, x_donut, y_donut, donut_width, donut_heigth, sentiment_counts, graph_color=graph_color)


                    # Stacked progress bar chart for sentiment by day
                    # Position stacked bar chart in bottom half, full width
                    x_bar = Inches(0.5)
                    y_bar = Inches(4.5)  # Start below top section
                    cx_bar = full_content_width
                    cy_bar = Inches(2.5)  # Remaining height

                    try:
                        bg_box = add_bg_box(slide10, x_bar, y_bar, cx_bar, cy_bar, color=CHART_BG_COLOR)

                        # Group data by Day instead of Company - with error handling
                        linkedin_data_copy = linkedin_data.copy()
                        linkedin_data_filtered = linkedin_data_copy[linkedin_data_copy['Sentiment'].isin([-1, 0, 1])]
                        
                        # Convert Day column to string to avoid datetime comparison issues
                        if 'Day' in linkedin_data_filtered.columns:
                            linkedin_data_filtered.loc[:, 'Day'] = linkedin_data_filtered['Day'].astype(str)
                        
                        day_sentiments = linkedin_data_filtered.groupby('Day')['Sentiment'].value_counts().unstack(fill_value=0)
                        
                        # Sort by day (now as strings)
                        try:
                            day_sentiments = day_sentiments.sort_index()
                        except Exception:
                            # If sorting fails, keep original order
                            pass
                        
                        chart_data = CategoryChartData()
                        chart_data.categories = [str(day) for day in day_sentiments.index.tolist()]
                        
                        # Add series with safe value extraction
                        for sentiment in [1, 0, -1]:
                            series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                            if sentiment in day_sentiments.columns:
                                try:
                                    series_values = [int(val) for val in day_sentiments[sentiment].tolist()]
                                    chart_data.add_series(series_name, series_values)
                                except Exception as e:
                                    print(f"Error adding series {series_name}: {e}")
                                    # Add empty series as fallback
                                    chart_data.add_series(series_name, [0] * len(day_sentiments))
                            else:
                                # Add empty series if sentiment not found
                                chart_data.add_series(series_name, [0] * len(day_sentiments))
                        
                        chart = slide10.shapes.add_chart(
                            XL_CHART_TYPE.COLUMN_STACKED,
                            x_bar, y_bar,
                            cx_bar, cy_bar,
                            chart_data
                        ).chart
                        
                        chart.has_legend = True
                        chart.has_data_labels = True
                        chart.legend.position = XL_LEGEND_POSITION.TOP
                        chart.legend.font.size = Pt(12)
                        
                        # Apply formatting and colors
                        try:
                            apply_chart_formatting(chart, title="Günlük sentiment dağılımı", graph_color=graph_color)
                            apply_sentiment_colors(chart)
                        except Exception as e:
                            print(f"Error applying chart formatting: {e}")
                        
                        # Add data labels at center with size 10
                        try:
                            for series in chart.series:
                                series.has_data_labels = True
                                data_labels = series.data_labels
                                data_labels.position = XL_DATA_LABEL_POSITION.CENTER
                                data_labels.font.size = Pt(10)
                                data_labels.font.bold = False
                                data_labels.font.color.rgb = RGBColor(255, 255, 255)  # White text for better visibility on colored background
                        except Exception as e:
                            print(f"Error setting data labels: {e}")
                            
                    except Exception as e:
                        print(f"Error creating stacked bar chart: {e}")
                else:
                    logger.warning(f"No Linkedin data found for company: {company_name}")
        # endregion

        # region eleveth slide - Linkedin sentiment analysis
        logger.debug("Creating eleveth slide with Linkedin sentiment analysis")
        slide11 = prs.slides.add_slide(prs.slide_layouts[5])
        # Remove default textbox
        for shape in slide11.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)
        add_slide_header(slide11, company_logo_path, start_date, end_date, "Linkedln postlarının analizi", title_color, template_color)
        add_side_line(slide11, template_color)

        # Set slide background
        background = slide11.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR
        
        if 'combined_sources' in data_frames and 'Linkedin' in data_frames['combined_sources']:
            linkedin_data = data_frames['combined_sources']['Linkedin']
            if has_competitors:
                company_sentiment = linkedin_data[linkedin_data['Company'] == company_name]
            else:
                company_sentiment = linkedin_data
            
            if not company_sentiment.empty:
                # Calculate heights accounting for header
                available_height = Inches(7.5 - HEADER_HEIGHT)  # Total height minus header
                half_height = available_height / 2
                full_content_width = Inches(12.33)
                if has_competitors:
                    # Top half section for donut and line charts
                    donut_size = Inches(3.5)
                    
                    # Donut chart on left top half
                    x_donut = Inches(0.5)
                    y_donut = Inches(1.2)  # Just below header
                    create_sentiment_donut_chart(slide11, x_donut, y_donut, donut_size, 
                        donut_size - Inches(0.4), sentiment_counts, graph_color=graph_color)
                                            
                    # Multiline chart in right half
                    # Position multiline chart in right half of top section
                    x_line = Inches(4.5)  # Start after donut chart
                    y_line = Inches(1.2)  # Same vertical alignment as donut
                    cx_line = Inches(8.33)  # Remaining width
                    cy_line = Inches(3.1)

                    sentiment_by_date = company_sentiment.groupby('Day')['Sentiment'].value_counts().unstack(fill_value=0)
                    chart_data = CategoryChartData()
                    chart_data.categories = sentiment_by_date.index.tolist()
                    
                    for sentiment in [1, 0, -1]:
                        series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                        if sentiment in sentiment_by_date.columns:
                            chart_data.add_series(series_name, sentiment_by_date[sentiment].tolist())
                    
                    create_sentiment_line_chart(slide11, x_line, y_line, cx_line, cy_line, chart_data, title="Postların sentiment və zamana görə bölgüsü", graph_color=graph_color)

                    # Vertical multibar chart for LinkedIn company sentiment comparison
                    # Position multibar chart in bottom half, full width
                    x_bar = Inches(0.5)
                    y_bar = Inches(4.5)  # Start below top section
                    cx_bar = Inches(12.33)  # Full width
                    cy_bar = Inches(2.5)  # Remaining height

                    bg_box = add_bg_box(slide11, x_bar, y_bar, cx_bar, cy_bar, color=CHART_BG_COLOR)

                    linkedin_data = linkedin_data[linkedin_data['Sentiment'].isin([-1, 0, 1])]
                    company_sentiments = linkedin_data.groupby('Company')['Sentiment'].value_counts().unstack(fill_value=0)
                    
                    # Sort by total sentiment values
                    totals = company_sentiments.sum(axis=1)
                    company_sentiments = company_sentiments.loc[totals.sort_values(ascending=False).index]
                    
                    chart_data = CategoryChartData()
                    chart_data.categories = company_sentiments.index.tolist()
                    
                    for sentiment in [1, 0, -1]:
                        series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                        if sentiment in company_sentiments.columns:
                            chart_data.add_series(series_name, company_sentiments[sentiment].tolist())
                    
                    chart = slide11.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        x_bar, y_bar,
                        cx_bar, cy_bar,
                        chart_data
                    ).chart
                    
                    chart.has_legend = True
                    chart.has_data_labels = True
                    chart.legend.position = XL_LEGEND_POSITION.TOP
                    chart.legend.font.size = Pt(12)
                    
                    # Apply formatting and colors
                    apply_chart_formatting(chart, title="Post saylarına görə bankların bölgüsü", graph_color=graph_color)
                    apply_sentiment_colors(chart)
                    
                    # Add data labels at outside end with size 10
                    for series in chart.series:
                        series.has_data_labels = True
                        data_labels = series.data_labels
                        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                        data_labels.font.size = Pt(10)
                        data_labels.font.bold = False
                        data_labels.font.color.rgb = RGBColor(89, 89, 89)
                else:  # has no cempoetitors
                    # Top half section for text and donut charts
                    donut_width = Inches(5)
                    donut_heigth = Inches(3.1)
                    
                    # Text section on left top half (where donut was)
                    x_text = Inches(0.5)
                    y_text = Inches(1.2)  # Just below header
                    cx_text = full_content_width - donut_width - Inches(0.5)  # Remaining width after donut
                    cy_text = Inches(3.1)  # Height for text section
                    
                    try:
                        # Add background box for text section
                        bg_box = add_bg_box(slide11, x_text, y_text, cx_text, cy_text, color=CHART_BG_COLOR)
                        
                        # Title text box
                        title_left = x_text + Inches(0.2)
                        title_top = y_text + Inches(0.2)
                        title_width = cx_text - Inches(0.4)
                        title_height = Inches(0.8)  # Fixed height for title
                        title_textbox = slide11.shapes.add_textbox(title_left, title_top, title_width, title_height)
                        title_tf = title_textbox.text_frame
                        title_tf.clear()
                        title_tf.word_wrap = True

                        # Title paragraph with colored text
                        title_p = title_tf.paragraphs[0]
                        title_p.clear()

                        # Add icon and styled title
                        run_icon = title_p.add_run()
                        run_icon.text = f"{CHARTS_ICONS['Time Distribution']} LinkedIn postlarında "
                        run_icon.font.size = Pt(16)
                        run_icon.font.bold = True
                        run_icon.font.color.rgb = RGBColor(0, 0, 0)

                        run_negative = title_p.add_run()
                        run_negative.text = "Neqativ "
                        run_negative.font.size = Pt(16)
                        run_negative.font.bold = True
                        run_negative.font.color.rgb = RGBColor(220, 53, 69)  # Red

                        run_text = title_p.add_run()
                        run_text.text = "xəbər qeydə alınmamışdır"
                        run_text.font.size = Pt(16)
                        run_text.font.bold = True
                        run_text.font.color.rgb = RGBColor(0, 0, 0)  # Black

                        run_positive = title_p.add_run()
                        run_positive.text = "Pozitiv "
                        run_positive.font.size = Pt(16)
                        run_positive.font.bold = True
                        run_positive.font.color.rgb = RGBColor(40, 167, 69)  # Green

                        run_text2 = title_p.add_run()
                        run_text2.text = "məzmunlu xəbərlər aşağıda qeyd edilmişdir."
                        run_text2.font.size = Pt(16)
                        run_text2.font.bold = True
                        run_text2.font.color.rgb = RGBColor(0, 0, 0)  # Black

                        title_p.alignment = PP_ALIGN.LEFT

                        # Second paragraph text box with bullet and colored text
                        para_left = x_text + Inches(0.2)
                        para_top = title_top + title_height + Inches(0.1)
                        para_width = cx_text - Inches(0.4)
                        para_height = Inches(1.5)  # Adjusted height for content
                        para_textbox = slide11.shapes.add_textbox(para_left, para_top, para_width, para_height)
                        para_tf = para_textbox.text_frame
                        para_tf.clear()
                        para_tf.word_wrap = True

                        para = para_tf.paragraphs[0]
                        para.level = 0  # Top-level bullet
                        para.clear()

                        # Add text with colored "neqativ" word
                        run1 = para.add_run()
                        run1.text = "4SİM Milli Proqramı İqtisadiyyat Nazirliyinin tabeliyində Dördüncü Sənaye İnqilabının Təhlili və Koordinasiya Mərkəzi, Elm və Təhsil Nazirliyinin tabeliyində Təhsilin İnkişafi Fondu və “Coursera” şirkətinin birgə əməkdaşlığı və “State Oil Company of the Republic of Azerbaijan”, ”PASHA Holding”, “bp” və “JOCAP” şirkətlərinin dəstəyi ilə icra olunur.və bu kimi xəbərlər "
                        run1.font.size = Pt(16)
                        run1.font.bold = False
                        run1.font.color.rgb = RGBColor(51, 51, 51)

                        run_positive = para.add_run()
                        run_positive.text = "positiv"
                        run_positive.font.size = Pt(16)
                        run_positive.font.bold = True
                        run_positive.font.color.rgb = RGBColor(40, 167, 69)  # Green

                        run2 = para.add_run()
                        run2.text = " olaraq qeyd edilmişdir."
                        run2.font.size = Pt(16)
                        run2.font.bold = False
                        run2.font.color.rgb = RGBColor(51, 51, 51)

                        para.alignment = PP_ALIGN.LEFT
                        
                    except Exception as e:
                        print(f"Error creating text section: {e}")
                    
                    # Donut chart on right top half (where multiline was)
                    x_donut = full_content_width - donut_width + Inches(0.5) # Start after text section
                    y_donut = Inches(1.2)  # Same vertical alignment as text
                    create_sentiment_donut_chart(slide11, x_donut, y_donut, donut_width, donut_heigth, sentiment_counts, graph_color=graph_color)

                    # Stacked progress bar chart for sentiment by day
                    # Position stacked bar chart in bottom half, full width
                    x_bar = Inches(0.5)
                    y_bar = Inches(4.5)  # Start below top section
                    cx_bar = full_content_width
                    cy_bar = Inches(2.5)  # Remaining height

                    try:
                        bg_box = add_bg_box(slide11, x_bar, y_bar, cx_bar, cy_bar, color=CHART_BG_COLOR)

                        # Group data by Day instead of Company - with error handling
                        linkedin_data_copy = linkedin_data.copy()
                        linkedin_data_filtered = linkedin_data_copy[linkedin_data_copy['Sentiment'].isin([-1, 0, 1])]
                        
                        # Convert Day column to string to avoid datetime comparison issues
                        if 'Day' in linkedin_data_filtered.columns:
                            linkedin_data_filtered.loc[:, 'Day'] = linkedin_data_filtered['Day'].astype(str)
                        
                        day_sentiments = linkedin_data_filtered.groupby('Day')['Sentiment'].value_counts().unstack(fill_value=0)
                        
                        # Sort by day (now as strings)
                        try:
                            day_sentiments = day_sentiments.sort_index()
                        except Exception:
                            # If sorting fails, keep original order
                            pass
                        
                        chart_data = CategoryChartData()
                        chart_data.categories = [str(day) for day in day_sentiments.index.tolist()]
                        
                        # Add series with safe value extraction
                        for sentiment in [1, 0, -1]:
                            series_name = "Positive" if sentiment == 1 else "Neutral" if sentiment == 0 else "Negative"
                            if sentiment in day_sentiments.columns:
                                try:
                                    series_values = [int(val) for val in day_sentiments[sentiment].tolist()]
                                    chart_data.add_series(series_name, series_values)
                                except Exception as e:
                                    print(f"Error adding series {series_name}: {e}")
                                    # Add empty series as fallback
                                    chart_data.add_series(series_name, [0] * len(day_sentiments))
                            else:
                                # Add empty series if sentiment not found
                                chart_data.add_series(series_name, [0] * len(day_sentiments))
                        
                        chart = slide11.shapes.add_chart(
                            XL_CHART_TYPE.COLUMN_STACKED,
                            x_bar, y_bar,
                            cx_bar, cy_bar,
                            chart_data
                        ).chart
                        
                        chart.has_legend = True
                        chart.has_data_labels = True
                        chart.legend.position = XL_LEGEND_POSITION.TOP
                        chart.legend.font.size = Pt(12)
                        
                        # Apply formatting and colors
                        try:
                            apply_chart_formatting(chart, title="Günlük sentiment dağılımı", graph_color=graph_color)
                            apply_sentiment_colors(chart)
                        except Exception as e:
                            print(f"Error applying chart formatting: {e}")
                        
                        # Add data labels at center with size 10
                        try:
                            for series in chart.series:
                                series.has_data_labels = True
                                data_labels = series.data_labels
                                data_labels.position = XL_DATA_LABEL_POSITION.CENTER
                                data_labels.font.size = Pt(10)
                                data_labels.font.bold = False
                                data_labels.font.color.rgb = RGBColor(255, 255, 255)  # White text for better visibility on colored background
                        except Exception as e:
                            print(f"Error setting data labels: {e}")
                            
                    except Exception as e:
                        print(f"Error creating stacked bar chart: {e}")
            else:
                logger.warning(f"No Linkedin data found for company: {company_name}")
        # endregion

        # region Twelveth slide - Positive and Negative Posts
        logger.debug("Creating Twelveth slide with positive and negative posts")
        slide12 = prs.slides.add_slide(prs.slide_layouts[5])
        # Remove default textbox
        for shape in slide12.shapes:
            if shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)
        add_slide_header(slide12, company_logo_path, start_date, end_date, "Sosial media postlarının analizi", title_color, template_color)
        add_side_line(slide12, template_color)

        # Layout constants
        header_height = Inches(0.8)
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)
        available_height = slide_height - header_height - Inches(0.4)  # Bottom margin
        content_top = header_height + Inches(0.2)

        # Set slide background
        background = slide12.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_COLOR

        # Post layout settings
        image_width = Inches(2.2)
        vertical_spacing = Inches(0.15)
        horizontal_spacing = Inches(0.3)  # Space between images in same row
        caption_height = Inches(0.3)
        group_top = content_top + Inches(0.8)  # Space for section titles
        max_group_height = available_height - Inches(1.2)  # Space for titles and clue card

        # Section titles
        title_height = Inches(0.4)
        title_top = content_top + Inches(0.2)

        # Negative posts title (left side)
        negative_title_left = Inches(0.5)
        negative_title_width = Inches(6)

        # Add white background shape for negative title
        negative_bg = slide12.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            negative_title_left, title_top, negative_title_width, title_height
        )
        negative_bg.fill.solid()
        negative_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        negative_bg.line.fill.background()
        negative_bg.line.width = Inches(0.01)

        # Add shadow to negative title background
        negative_bg.shadow.inherit = True
        negative_bg.shadow.blur_radius = 3000
        negative_bg.shadow.distance = 2000
        negative_bg.shadow.angle = 45

        negative_title_box = slide12.shapes.add_textbox(
            negative_title_left, title_top, negative_title_width, title_height
        )
        negative_tf = negative_title_box.text_frame
        negative_tf.clear()
        negative_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        negative_p = negative_tf.paragraphs[0]

        # Create runs for different styling
        negative_run1 = negative_p.runs[0] if negative_p.runs else negative_p.add_run()
        negative_run1.text = "Neqativ "
        negative_run1.font.size = Pt(18)
        negative_run1.font.bold = True
        negative_run1.font.color.rgb = RGBColor(220, 53, 69)  # Red color

        negative_run2 = negative_p.add_run()
        negative_run2.text = "post nümunələri"
        negative_run2.font.size = Pt(18)
        negative_run2.font.bold = False
        negative_run2.font.color.rgb = RGBColor(0, 0, 0)  # Black color

        negative_p.alignment = PP_ALIGN.CENTER

        # Positive posts title (right side)
        positive_title_left = Inches(6.83)  # Right half starts here
        positive_title_width = Inches(6)

        # Add white background shape for positive title
        positive_bg = slide12.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            positive_title_left, title_top, positive_title_width, title_height
        )
        positive_bg.fill.solid()
        positive_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        positive_bg.line.fill.background()
        positive_bg.line.width = Inches(0.01)

        # Add shadow to positive title background
        positive_bg.shadow.inherit = True
        positive_bg.shadow.blur_radius = 3000
        positive_bg.shadow.distance = 2000
        positive_bg.shadow.angle = 45

        positive_title_box = slide12.shapes.add_textbox(
            positive_title_left, title_top, positive_title_width, title_height
        )
        positive_tf = positive_title_box.text_frame
        positive_tf.clear()
        positive_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        positive_p = positive_tf.paragraphs[0]

        # Create runs for different styling
        positive_run1 = positive_p.runs[0] if positive_p.runs else positive_p.add_run()
        positive_run1.text = "Pozitiv "
        positive_run1.font.size = Pt(18)
        positive_run1.font.bold = True
        positive_run1.font.color.rgb = RGBColor(40, 167, 69)  # Green color

        positive_run2 = positive_p.add_run()
        positive_run2.text = "post nümunələri"
        positive_run2.font.size = Pt(18)
        positive_run2.font.bold = False
        positive_run2.font.color.rgb = RGBColor(0, 0, 0)  # Black color

        positive_p.alignment = PP_ALIGN.CENTER

        def add_post(slide, post, left, top):
            if not os.path.exists(post["image_path"]):
                return 0

            # Add image and get actual height
            img = slide.shapes.add_picture(post["image_path"], left, top, width=image_width)
            img_height = img.height / 914400

            # Link
            img.click_action.hyperlink.address = post["link"]

            # Styling
            img.line.color.rgb = RGBColor(200, 200, 200)
            img.line.width = Inches(0.02)
            img.shadow.inherit = False
            img.shadow.blur_radius = 5000

            return Inches(img_height) + caption_height + vertical_spacing

        def layout_posts(slide, posts, group_center_x):
            count = len(posts)
            positions = []
            
            if count == 0:
                return

            if count == 1:
                positions.append((group_center_x - image_width / 2, group_top + Inches(1.2)))
            elif count == 2:
                positions.append((group_center_x - image_width / 2, group_top + Inches(0.2)))  # top
                positions.append((group_center_x - image_width / 2, group_top + Inches(2.2)))  # bottom with more space
            elif count >= 3:
                # Top row (two side by side) - Adjusted spacing for bigger images
                left1 = group_center_x - image_width - horizontal_spacing / 2
                left2 = group_center_x + horizontal_spacing / 2
                top1 = group_top + Inches(0.1)
                positions.append((left1, top1))
                positions.append((left2, top1))
                # Centered below with more space
                center_x = group_center_x - image_width / 2
                positions.append((center_x, group_top + Inches(2.1)))

            # Add posts (limit to 3)
            for post, (left, top) in zip(posts[:3], positions):
                add_post(slide, post, left, top)

        # Apply layout for each group
        if negative_posts:
            layout_posts(slide12, negative_posts, group_center_x=Inches(3.25))  # Left half center

        if positive_posts:
            layout_posts(slide12, positive_posts, group_center_x=Inches(9.75))  # Right half center

        clue_card_width = Inches(2)
        clue_card_height = Inches(1.5)
        clue_card_left = (slide_width - clue_card_width) / 2  # Center horizontally
        clue_card_top = slide_height - clue_card_height - Inches(0.2)  # Bottom with margin

        # Red background card with shadow
        clue_bg = slide12.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            clue_card_left, clue_card_top, clue_card_width, clue_card_height
        )
        clue_bg.fill.solid()
        clue_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        clue_bg.line.color.rgb = RGBColor(255, 255, 255)
        clue_bg.line.width = Inches(0.01)

        # Add shadow to clue card
        clue_bg.shadow.inherit = False
        clue_bg.shadow.blur_radius = 4000
        clue_bg.shadow.distance = 2500
        clue_bg.shadow.angle = 45

        # Add stick icon (positioned at top-right corner, half outside)
        stick_size = Inches(0.4)
        stick_left = clue_card_left + clue_card_width - stick_size / 2  # Half outside right edge
        stick_top = clue_card_top - stick_size / 2  # Half outside top edge

        # Add stick icon text (using a pin/stick emoji or symbol)
        stick_text = slide12.shapes.add_textbox(
            stick_left, stick_top, stick_size, stick_size
        )
        stick_tf = stick_text.text_frame
        stick_tf.clear()
        stick_tf.margin_top = Inches(0)
        stick_tf.margin_bottom = Inches(0)
        stick_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        stick_p = stick_tf.paragraphs[0]
        stick_p.text = "📌"  # Pin/stick emoji
        stick_p.font.size = Pt(24)
        stick_p.alignment = PP_ALIGN.CENTER
        stick_p.font.color.rgb = graph_color

        # Clue text (adjusted for better positioning)
        clue_text = slide12.shapes.add_textbox(
            clue_card_left + Inches(0.1), clue_card_top, clue_card_width - Inches(0.2), clue_card_height
        )
        clue_tf = clue_text.text_frame
        clue_tf.clear()
        clue_tf.margin_top = Inches(0)
        clue_tf.margin_bottom = Inches(0)
        clue_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        clue_tf.word_wrap = True
        clue_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        clue_p = clue_tf.paragraphs[0]
        clue_p.text = "Şəkillərə klikləyərək orijinal postları görə bilərsiniz"
        clue_p.font.size = Pt(16)
        clue_p.font.color.rgb = graph_color
        clue_p.alignment = PP_ALIGN.CENTER

        # endregion

        logger.debug("Saving PowerPoint file")
        prs.save(output_path)
        logger.debug("PowerPoint file saved successfully")
    except Exception as e:
        logger.error(f"Error creating PowerPoint: {str(e)}")
        raise